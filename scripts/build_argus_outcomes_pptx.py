#!/usr/bin/env python3
"""Generate a clean PowerPoint deck from docs/argus_outcomes.md.

デザイン方針:
- テーマ: Ocean Gradient (deep navy / teal / ice white) + warm accent(coral)
- 素人向け: 各コマンドに「ひとことで言うと」メタファー + 絵文字アイコン
- モチーフ: 番号バッジ (teal円) + 左サイドの navy バー を全コンテンツスライドで反復
- サンドイッチ構造: 表紙/まとめは dark、コンテンツは light
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(__file__).resolve().parent.parent / "docs" / "argus_outcomes.pptx"

# === Ocean Gradient palette ===
NAVY   = RGBColor(0x0B, 0x2A, 0x4A)   # dominant
DEEP   = RGBColor(0x06, 0x5A, 0x82)
TEAL   = RGBColor(0x1C, 0x72, 0x93)   # secondary
MINT   = RGBColor(0x9A, 0xD1, 0xD4)
ICE    = RGBColor(0xEE, 0xF5, 0xF9)   # light bg
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x16, 0x1F, 0x2C)
GRAY   = RGBColor(0x5A, 0x6A, 0x7E)
MUTED  = RGBColor(0x9A, 0xA7, 0xB8)
CORAL  = RGBColor(0xF1, 0x6B, 0x4F)   # warm accent
GOLD   = RGBColor(0xE8, 0xB5, 0x3B)
SUCCESS= RGBColor(0x4C, 0xAF, 0x50)

HEADER_FONT = "Meiryo"
BODY_FONT   = "Meiryo"
MONO_FONT   = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- primitives ----------
def add_rect(slide, x, y, w, h, color, *, line=None, rounded=False, radius=0.12):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line; shape.line.width = Pt(0.75)
    if rounded:
        shape.adjustments[0] = radius
    return shape


def add_bg(slide, color=WHITE):
    return add_rect(slide, 0, 0, SW, SH, color)


def add_text(slide, x, y, w, h, text, *, size=16, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font=BODY_FONT, anchor=None, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=DARK, bullet_color=None, gap=4):
    bullet_color = bullet_color or TEAL
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_top = Inches(0.02)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        r1 = p.add_run(); r1.text = "▸ "
        r1.font.name = BODY_FONT; r1.font.size = Pt(size); r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run(); r2.text = it
        r2.font.name = BODY_FONT; r2.font.size = Pt(size); r2.font.color.rgb = color
    return tb


def code_block(slide, x, y, w, h, lines, *, size=12):
    add_rect(slide, x, y, w, h, RGBColor(0x10, 0x20, 0x30), rounded=True, radius=0.05)
    # accent dot strip (macOS-ish window chrome)
    for i, c in enumerate([CORAL, GOLD, SUCCESS]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.18 + i*0.22), y + Inches(0.15),
                                     Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = c
        dot.line.fill.background()
    tb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.4),
                                  w - Inches(0.5), h - Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run(); r.text = ln
        r.font.name = MONO_FONT; r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0xDA, 0xE8, 0xF4)


def number_badge(slide, x, y, num, *, size_in=0.7, bg=TEAL, fg=WHITE):
    d = Inches(size_in)
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    c.fill.solid(); c.fill.fore_color.rgb = bg
    c.line.fill.background()
    add_text(slide, x, y, d, d, str(num), size=int(size_in * 28),
             bold=True, color=fg, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font=HEADER_FONT)


def header(slide, number, title, subtitle=None):
    # left vertical navy bar (motif)
    add_rect(slide, 0, 0, Inches(0.35), SH, NAVY)
    add_rect(slide, Inches(0.35), 0, Inches(0.06), SH, TEAL)
    # top title area
    add_rect(slide, Inches(0.55), Inches(0.45), Inches(12.5), Inches(0.05), MINT)
    number_badge(slide, Inches(0.65), Inches(0.65), number, size_in=0.85)
    add_text(slide, Inches(1.75), Inches(0.65), Inches(11.2), Inches(0.55),
             title, size=24, bold=True, color=NAVY, font=HEADER_FONT)
    if subtitle:
        add_text(slide, Inches(1.75), Inches(1.20), Inches(11.2), Inches(0.4),
                 subtitle, size=13, color=GRAY, italic=True)
    # separator under header
    add_rect(slide, Inches(0.55), Inches(1.72), Inches(12.5), Inches(0.02),
             RGBColor(0xE0, 0xE8, 0xF0))


def section_title(slide, x, y, w, text, *, icon=None):
    # Small vertical accent bar (motif) instead of emoji icon for font safety
    add_rect(slide, x, y + Inches(0.05), Inches(0.08), Inches(0.32), CORAL)
    x2 = x + Inches(0.2)
    w2 = w - Inches(0.2)
    add_text(slide, x2, y, w2, Inches(0.4), text, size=14, bold=True,
             color=TEAL, font=HEADER_FONT)


def callout(slide, x, y, w, h, label, text, *, bg=RGBColor(0xFF, 0xF6, 0xE5),
            border=GOLD, label_color=RGBColor(0x8A, 0x5A, 0x00),
            text_size=12):
    add_rect(slide, x, y, w, h, bg, line=border, rounded=True, radius=0.06)
    add_text(slide, x + Inches(0.2), y + Inches(0.12), w - Inches(0.4), Inches(0.35),
             label, size=12, bold=True, color=label_color)
    add_text(slide, x + Inches(0.2), y + Inches(0.48), w - Inches(0.4), h - Inches(0.55),
             text, size=text_size, color=DARK)


def footer(slide, text="Argus AI — 使い方ガイド", page=None, total=None):
    add_text(slide, Inches(0.55), SH - Inches(0.35), Inches(8), Inches(0.25),
             text, size=9, color=MUTED)
    if page is not None:
        add_text(slide, SW - Inches(1.2), SH - Inches(0.35), Inches(0.8), Inches(0.25),
                 f"{page} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ========== Slide 1: Cover ==========
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
# gradient effect via overlayed bands
add_rect(s, 0, Inches(5.8), SW, Inches(1.7), DEEP)
add_rect(s, 0, Inches(7.3), SW, Inches(0.2), TEAL)
# decorative circle
c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1.5), Inches(5.5), Inches(5.5))
c.fill.solid(); c.fill.fore_color.rgb = DEEP
c.line.fill.background()
c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.2), Inches(-0.8), Inches(4.0), Inches(4.0))
c2.fill.solid(); c2.fill.fore_color.rgb = TEAL
c2.line.fill.background()

add_text(s, Inches(0.8), Inches(1.6), Inches(3), Inches(0.5),
         "USER GUIDE", size=14, bold=True, color=MINT, font=HEADER_FONT)
add_rect(s, Inches(0.8), Inches(2.05), Inches(0.8), Inches(0.08), CORAL)

add_text(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.6),
         "Argus AI", size=88, bold=True, color=WHITE, font=HEADER_FONT)
add_text(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.7),
         "プロジェクトマネージャーのための\nSlack コマンド 5 選",
         size=26, color=MINT)
add_text(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.5),
         "状況把握・リスク検知・議事録化 ― すべて Slack から一声で。",
         size=14, color=RGBColor(0xB8, 0xCC, 0xE0), italic=True)
add_text(s, Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.3),
         "docs/argus_outcomes.md より", size=10, color=MUTED)


# ========== Slide 2: ひとことで言うと (metaphor overview) ==========
s = prs.slides.add_slide(BLANK)
add_bg(s, ICE)
header(s, "0", "ひとことで言うと", "各コマンドを日常の役割にたとえると…")

metaphors = [
    ("📋", "/argus-brief",      "朝の秘書ブリーフィング",
     "「今日、何から手をつけるべき?」を5件にまとめてくれる。"),
    ("🌇", "/argus-today",      "退勤前のラップアップ",
     "今日起きたこと + 自分宛の依頼を見落とさずに受け取る。"),
    ("🚨", "/argus-risk",       "火種を見つける番人",
     "顕在リスクと予兆を優先度で並べ、役員会前に慌てない。"),
    ("🔎", "/argus-investigate","調査員 (探偵)",
     "「なぜ遅れた?」「どこで決めた?」を自律的に追跡・回答。"),
    ("🎙️","/argus-transcribe", "プロの速記者",
     "会議録音を投げれば、議事録と決定事項を自動で書き起こす。"),
]
col_w = Inches(2.48); gap = Inches(0.1); left = Inches(0.55)
top = Inches(2.05); card_h = Inches(4.6)
for i, (icon, cmd, role, desc) in enumerate(metaphors):
    x = left + (col_w + gap) * i
    # card
    add_rect(s, x, top, col_w, card_h, WHITE,
             line=RGBColor(0xD6, 0xE0, 0xEA), rounded=True, radius=0.04)
    # colored top strip
    stripe_colors = [TEAL, DEEP, CORAL, NAVY, SUCCESS]
    add_rect(s, x, top, col_w, Inches(0.12), stripe_colors[i])
    # icon in circle
    cx = x + col_w/2 - Inches(0.55)
    cy = top + Inches(0.4)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(1.1), Inches(1.1))
    circ.fill.solid(); circ.fill.fore_color.rgb = ICE
    circ.line.color.rgb = stripe_colors[i]; circ.line.width = Pt(1.5)
    add_text(s, cx, cy, Inches(1.1), Inches(1.1), icon,
             size=44, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # command chip
    add_text(s, x + Inches(0.1), top + Inches(1.75), col_w - Inches(0.2), Inches(0.4),
             cmd, size=13, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, font=MONO_FONT)
    # role
    add_text(s, x + Inches(0.15), top + Inches(2.3), col_w - Inches(0.3), Inches(0.5),
             role, size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    # desc
    add_text(s, x + Inches(0.2), top + Inches(3.0), col_w - Inches(0.4), Inches(1.5),
             desc, size=11, color=GRAY, align=PP_ALIGN.CENTER)

footer(s, page=2, total=8)


# ========== Slide 3: Overview (detail table) ==========
s = prs.slides.add_slide(BLANK)
add_bg(s, ICE)
header(s, "1", "5 コマンドの全体像",
       "すべて ephemeral 返信（あなただけに見える） / transcribe の進捗のみスレッド公開")

commands = [
    ("📋", "/argus-brief",      "優先順位づけブリーフィング", "今週〜直近を俯瞰して優先度付け",     "30–60秒",  TEAL),
    ("🌇", "/argus-today",      "今日の活動サマリー（個人）", "本日の4観点 + 自分宛メンション",       "20–40秒",  DEEP),
    ("🚨", "/argus-risk",       "リスク分析",                "顕在リスクと予兆を優先度付きで列挙", "30–60秒",  CORAL),
    ("🔎", "/argus-investigate","マルチステップ調査（Agent）","「なぜ」「どこで決めた」の深掘り",     "30–180秒", NAVY),
    ("🎙️","/argus-transcribe", "会議録音 → 議事録生成",    "Whisper + LLM で自動議事録化",        "10–20分",  SUCCESS),
]

top = Inches(2.0)
row_h = Inches(0.95); row_gap = Inches(0.08)
for i, (icon, cmd, name, desc, dur, clr) in enumerate(commands):
    y = top + (row_h + row_gap) * i
    add_rect(s, Inches(0.55), y, Inches(12.3), row_h, WHITE,
             line=RGBColor(0xD6, 0xE0, 0xEA), rounded=True, radius=0.15)
    # left accent bar
    add_rect(s, Inches(0.55), y, Inches(0.15), row_h, clr, rounded=True, radius=0.3)
    # icon
    add_text(s, Inches(0.85), y, Inches(0.9), row_h, icon,
             size=28, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # command chip
    chip_w = Inches(2.8)
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(1.85), y + Inches(0.25), chip_w, Inches(0.45))
    chip.fill.solid(); chip.fill.fore_color.rgb = NAVY
    chip.line.fill.background(); chip.adjustments[0] = 0.3
    add_text(s, Inches(1.85), y + Inches(0.28), chip_w, Inches(0.4), cmd,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=MONO_FONT)
    # name
    add_text(s, Inches(4.8), y + Inches(0.08), Inches(5.7), Inches(0.42),
             name, size=14, bold=True, color=DARK)
    # desc
    add_text(s, Inches(4.8), y + Inches(0.48), Inches(5.7), Inches(0.42),
             desc, size=11, color=GRAY)
    # duration pill
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(10.85), y + Inches(0.28), Inches(1.85), Inches(0.4))
    pill.fill.solid(); pill.fill.fore_color.rgb = ICE
    pill.line.color.rgb = clr; pill.line.width = Pt(1.0); pill.adjustments[0] = 0.5
    add_text(s, Inches(10.85), y + Inches(0.3), Inches(1.85), Inches(0.36), "⏱ " + dur,
             size=11, bold=True, color=clr, align=PP_ALIGN.CENTER)

footer(s, page=3, total=8)


# ========== common content slide helper ==========
def content_slide(num, title, subtitle, metaphor_icon, metaphor_text):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, ICE)
    header(s, num, title, subtitle)
    # "ひとことで言うと" top banner (metaphor, beginner-friendly)
    bx, by, bw, bh = Inches(0.55), Inches(1.85), Inches(12.5), Inches(0.75)
    add_rect(s, bx, by, bw, bh, WHITE, line=MINT, rounded=True, radius=0.15)
    add_rect(s, bx, by, Inches(0.15), bh, CORAL, rounded=True, radius=0.5)
    add_text(s, bx + Inches(0.3), by, Inches(0.8), bh, metaphor_icon,
             size=28, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, bx + Inches(1.2), by + Inches(0.08), Inches(2.2), Inches(0.3),
             "ひとことで言うと", size=10, bold=True, color=CORAL)
    add_text(s, bx + Inches(1.2), by + Inches(0.3), bw - Inches(1.3), Inches(0.45),
             metaphor_text, size=14, bold=True, color=NAVY)
    return s


# ========== Slide 4: /argus-brief ==========
s = content_slide("2", "/argus-brief — 優先順位づけブリーフィング",
                  "pm.db 統計 × Slack × 議事録 を総合分析し、優先度順に最大5件",
                  "📋", "朝イチに「今日なにから着手?」を秘書が5件にまとめてくれる。")

section_title(s, Inches(0.55), Inches(2.85), Inches(6), "使い方", icon="💬")
code_block(s, Inches(0.55), Inches(3.3), Inches(6.3), Inches(2.1), [
    "/argus-brief",
    "/argus-brief 60",
    "/argus-brief @富岳太郎",
    "/argus-brief Benchpark",
    "/argus-brief 60 @富岳太郎 GPU性能",
])

section_title(s, Inches(0.55), Inches(5.55), Inches(6), "引数ルール", icon="🔤")
add_bullets(s, Inches(0.65), Inches(5.95), Inches(6.2), Inches(1.5), [
    "数字のみ  → 直近日数（例: 60 = 過去60日）",
    "@で始まる → 担当者フォーカス（例: @富岳太郎）",
    "その他文字 → 話題フォーカス（例: Benchpark）",
], size=12)

section_title(s, Inches(7.1), Inches(2.85), Inches(6), "こんなときに使う", icon="🎯")
scenes = [
    ("☕", "毎朝5分のレビュー", "出社して最初のコーヒーの間に"),
    ("📅", "週初めの山の把握", "「今週の重め」をチームで共有"),
    ("🧭", "マイルストーン確認", "遅延気味のMSで関連AI・決定を一望"),
]
sy = Inches(3.35)
for i, (emoji, t, d) in enumerate(scenes):
    y = sy + Inches(1.25) * i
    add_rect(s, Inches(7.1), y, Inches(5.9), Inches(1.1), WHITE,
             line=RGBColor(0xD6, 0xE0, 0xEA), rounded=True, radius=0.1)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(7.3), y + Inches(0.2), Inches(0.7), Inches(0.7))
    circ.fill.solid(); circ.fill.fore_color.rgb = MINT
    circ.line.fill.background()
    add_text(s, Inches(7.3), y + Inches(0.2), Inches(0.7), Inches(0.7), emoji,
             size=22, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.2), y + Inches(0.2), Inches(4.7), Inches(0.4),
             t, size=14, bold=True, color=DARK)
    add_text(s, Inches(8.2), y + Inches(0.58), Inches(4.7), Inches(0.5),
             d, size=11, color=GRAY)

footer(s, page=4, total=8)


# ========== Slide 5: /argus-today ==========
s = content_slide("3", "/argus-today — 今日の活動サマリー（個人向け）",
                  "本日の4観点（議論・決定・AI・進捗）＋ 自分宛メンション別セクション",
                  "🌇", "退勤前に「今日の総まとめ + 自分宛の依頼」を1コマンドで確認。")

section_title(s, Inches(0.55), Inches(2.85), Inches(6), "使い方", icon="💬")
code_block(s, Inches(0.55), Inches(3.3), Inches(6.3), Inches(1.0),
           ["/argus-today    # 引数なし・本日分のみ"])

# Timeline diagram (16:50 cron → 17:00 run)
section_title(s, Inches(0.55), Inches(4.5), Inches(6), "使うタイミング", icon="⏰")
tx, ty, tw = Inches(0.7), Inches(5.0), Inches(6.1)
# line
add_rect(s, tx, ty + Inches(0.35), tw, Inches(0.04), MUTED)
# marker 16:50
m1 = s.shapes.add_shape(MSO_SHAPE.OVAL, tx + Inches(0.8), ty + Inches(0.22),
                         Inches(0.3), Inches(0.3))
m1.fill.solid(); m1.fill.fore_color.rgb = TEAL; m1.line.fill.background()
add_text(s, tx + Inches(0.3), ty - Inches(0.1), Inches(1.5), Inches(0.3),
         "16:50", size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s, tx + Inches(0.1), ty + Inches(0.6), Inches(2.0), Inches(0.4),
         "cron で Slack 取込", size=10, color=GRAY, align=PP_ALIGN.CENTER)
# marker 17:00 (recommended)
m2 = s.shapes.add_shape(MSO_SHAPE.OVAL, tx + Inches(3.8), ty + Inches(0.15),
                         Inches(0.45), Inches(0.45))
m2.fill.solid(); m2.fill.fore_color.rgb = CORAL; m2.line.color.rgb = WHITE
m2.line.width = Pt(2)
add_text(s, tx + Inches(3.8), ty + Inches(0.18), Inches(0.45), Inches(0.4),
         "✓", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, tx + Inches(3.3), ty - Inches(0.1), Inches(1.5), Inches(0.3),
         "17:00〜", size=12, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
add_text(s, tx + Inches(3.0), ty + Inches(0.6), Inches(2.0), Inches(0.4),
         "ここで実行 (推奨)", size=10, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
add_text(s, tx + Inches(0.5), ty + Inches(1.1), tw, Inches(0.35),
         "17:00 より前に実行すると当日データが不完全になります。",
         size=10, color=GRAY, italic=True)

# comparison table (brief vs today)
section_title(s, Inches(7.1), Inches(2.85), Inches(6), "/argus-brief との違い", icon="⚖️")
headers_row = ["観点", "/argus-brief", "/argus-today"]
rows = [
    ["対象期間",   "直近30日",       "本日のみ"],
    ["出力形式",   "優先アクション5件","4観点サマリー"],
    ["メンション", "なし",           "自分宛を別枠"],
]
tbl_x, tbl_y = Inches(7.1), Inches(3.3)
col_w = [Inches(1.7), Inches(2.1), Inches(2.1)]
rh = Inches(0.5)
x = tbl_x
for i, htxt in enumerate(headers_row):
    add_rect(s, x, tbl_y, col_w[i], rh, NAVY, line=WHITE)
    add_text(s, x, tbl_y, col_w[i], rh, htxt,
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += col_w[i]
for ri, row in enumerate(rows):
    x = tbl_x; y = tbl_y + rh * (ri + 1)
    fc = ICE if ri % 2 == 0 else WHITE
    for ci, v in enumerate(row):
        add_rect(s, x, y, col_w[ci], rh, fc, line=RGBColor(0xD6, 0xE0, 0xEA))
        add_text(s, x, y, col_w[ci], rh, v,
                 size=11, color=DARK, align=PP_ALIGN.CENTER,
                 bold=(ci == 0), anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[ci]

# scene bullets
section_title(s, Inches(7.1), Inches(5.5), Inches(6), "こんなときに", icon="🎯")
add_bullets(s, Inches(7.2), Inches(5.9), Inches(5.8), Inches(1.5), [
    "17時以降のラップアップに",
    "依頼・決定・AI の取りこぼし防止",
    "個人視点での日次まとめ",
], size=12)

footer(s, page=5, total=8)


# ========== Slide 6: /argus-risk ==========
s = content_slide("4", "/argus-risk — リスク分析",
                  "顕在化しているリスクと、放置すると問題になる予兆を優先度付きで列挙",
                  "🚨", "プロジェクトの「火種」を H / M / L に並べてくれる番人。")

section_title(s, Inches(0.55), Inches(2.85), Inches(6), "使い方 (brief と同じ引数ルール)", icon="💬")
code_block(s, Inches(0.55), Inches(3.3), Inches(6.3), Inches(2.1), [
    "/argus-risk",
    "/argus-risk 60",
    "/argus-risk @富岳太郎",
    "/argus-risk Benchpark",
])

# Priority pyramid
section_title(s, Inches(0.55), Inches(5.55), Inches(6), "優先度の見方", icon="📊")
py = Inches(5.95)
levels = [
    ("H", "今週中に対応", CORAL),
    ("M", "計画に組み込む", GOLD),
    ("L", "観察継続",      SUCCESS),
]
for i, (lv, desc, clr) in enumerate(levels):
    y = py + Inches(0.4) * i
    add_rect(s, Inches(0.65), y, Inches(0.55), Inches(0.35), clr, rounded=True, radius=0.3)
    add_text(s, Inches(0.65), y, Inches(0.55), Inches(0.35), lv,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.35), y + Inches(0.03), Inches(5.2), Inches(0.3),
             desc, size=12, color=DARK)

# scenes cards
section_title(s, Inches(7.1), Inches(2.85), Inches(6), "こんなときに使う", icon="🎯")
scenes = [
    ("🏛️", "役員会前",      "「報告すべき火種」を洗い出す"),
    ("🔄", "責任範囲の引き継ぎ直後", "全体のリスク状況を俯瞰"),
    ("🔬", "brief の深掘り", "気になった項目をリスク観点で再評価"),
]
sy = Inches(3.35)
for i, (emoji, t, d) in enumerate(scenes):
    y = sy + Inches(1.25) * i
    add_rect(s, Inches(7.1), y, Inches(5.9), Inches(1.1), WHITE,
             line=RGBColor(0xD6, 0xE0, 0xEA), rounded=True, radius=0.1)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(7.3), y + Inches(0.2), Inches(0.7), Inches(0.7))
    circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0xFD, 0xE3, 0xDA)
    circ.line.fill.background()
    add_text(s, Inches(7.3), y + Inches(0.2), Inches(0.7), Inches(0.7), emoji,
             size=22, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.2), y + Inches(0.2), Inches(4.7), Inches(0.4),
             t, size=14, bold=True, color=DARK)
    add_text(s, Inches(8.2), y + Inches(0.58), Inches(4.7), Inches(0.5),
             d, size=11, color=GRAY)

footer(s, page=6, total=8)


# ========== Slide 7: /argus-investigate ==========
s = content_slide("5", "/argus-investigate — マルチステップ調査 (Agent)",
                  "LLM が DB / FTS / Slack のツールを自律選択、最大5ステップで深掘り",
                  "🔎", "探偵のように「なぜ?」「どこで?」を自分でツールを選んで追跡。")

section_title(s, Inches(0.55), Inches(2.85), Inches(6), "使い方", icon="💬")
code_block(s, Inches(0.55), Inches(3.3), Inches(6.3), Inches(2.3), [
    "/argus-investigate M3の遅延原因",
    "/argus-investigate 先週の決定の実行状況",
    "/argus-investigate @富岳太郎 の負荷が高い原因",
    "/argus-investigate 設計方針に関する最近の議論",
    "/argus-investigate GPU契約の最新版の場所",
])

section_title(s, Inches(0.55), Inches(5.75), Inches(6), "得意な質問", icon="🎯")
add_bullets(s, Inches(0.65), Inches(6.15), Inches(6.2), Inches(1.3), [
    "因果分析・クロスソース相関",
    "過去決定の検索・文書探索",
    "構造化QA (担当・件数など)",
], size=11, gap=2)

# Right: HyDE callout (the key "売り")
section_title(s, Inches(7.1), Inches(2.85), Inches(6), "HyDEクエリ拡張 — 取りこぼし防止")

# before/after diagram
box_x = Inches(7.1); box_y = Inches(3.3); box_w = Inches(5.9)
add_rect(s, box_x, box_y, box_w, Inches(4.0), WHITE,
         line=MINT, rounded=True, radius=0.04)

# 問題
add_text(s, box_x + Inches(0.2), box_y + Inches(0.15), box_w - Inches(0.4), Inches(0.3),
         "❌ 従来の全文検索の弱点", size=11, bold=True, color=CORAL)
add_text(s, box_x + Inches(0.3), box_y + Inches(0.48), box_w - Inches(0.5), Inches(0.7),
         "質問「配布を希望する人のリストは?」では、\n「富岳太郎 / 富岳花子…」と名前だけのチャンクをヒットできない。",
         size=10, color=GRAY)

# 解決
add_text(s, box_x + Inches(0.2), box_y + Inches(1.3), box_w - Inches(0.4), Inches(0.3),
         "✅ HyDE による解決", size=11, bold=True, color=SUCCESS)
# flow boxes
fx = box_x + Inches(0.3); fy = box_y + Inches(1.75); fw = Inches(1.55); fh = Inches(0.55)
gaps = Inches(0.2)
steps_flow = [("元の質問", TEAL), ("別表現 ×2\n(LLM生成)", DEEP), ("3クエリで\n並列検索", CORAL)]
for i, (t, clr) in enumerate(steps_flow):
    x = fx + (fw + gaps) * i
    add_rect(s, x, fy, fw, fh, clr, rounded=True, radius=0.15)
    add_text(s, x, fy, fw, fh, t, size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        arrow_x = x + fw + Inches(0.02)
        add_text(s, arrow_x, fy + Inches(0.1), Inches(0.18), Inches(0.4),
                 "→", size=16, bold=True, color=GRAY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# next row
fy2 = fy + Inches(0.8)
add_text(s, fx + Inches(0.5), fy2 + Inches(0.05), fw*3, Inches(0.3),
         "↓", size=14, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
fy3 = fy2 + Inches(0.35)
for i, (t, clr) in enumerate([("重複排除→マージ", MINT), ("元クエリで\nre-rank", GOLD), ("回答生成", SUCCESS)]):
    x = fx + (fw + gaps) * i
    add_rect(s, x, fy3, fw, fh, clr, rounded=True, radius=0.15)
    add_text(s, x, fy3, fw, fh, t, size=10, bold=True, color=DARK if i!=2 else WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        arrow_x = x + fw + Inches(0.02)
        add_text(s, arrow_x, fy3 + Inches(0.1), Inches(0.18), Inches(0.4),
                 "→", size=16, bold=True, color=GRAY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# applies to
add_text(s, box_x + Inches(0.2), box_y + Inches(3.7), box_w - Inches(0.4), Inches(0.3),
         "※ brief / risk / エンリッチメントにも同じ仕組みを適用",
         size=9, color=GRAY, italic=True)

footer(s, page=7, total=8)


# ========== Slide 8: /argus-transcribe ==========
s = content_slide("6", "/argus-transcribe — 会議録音 → 議事録生成",
                  "Whisper + スライドOCR + LLM でマルチステージ議事録化、スレッドに投稿",
                  "🎙️", "録音ファイルを投げると、プロの速記者が議事録・決定・AIまで抽出。")

section_title(s, Inches(0.55), Inches(2.85), Inches(6), "使い方", icon="💬")
code_block(s, Inches(0.55), Inches(3.3), Inches(6.3), Inches(1.2), [
    "/argus-transcribe GMT20260302-032528_Recording.mp4",
    "/argus-transcribe 2026-04-20_Leader_Meeting.m4a",
])

# Pipeline (horizontal flow)
section_title(s, Inches(0.55), Inches(4.65), Inches(6), "処理フロー", icon="⚙️")
pipe = [
    ("📥", "DL",       TEAL),
    ("🖼️", "OCR",    DEEP),
    ("👂", "Whisper", CORAL),
    ("📝", "議事録",   NAVY),
    ("📤", "投稿",     SUCCESS),
]
pfx, pfy = Inches(0.55), Inches(5.1)
box_w = Inches(1.08); gap = Inches(0.15)
for i, (ic, lab, clr) in enumerate(pipe):
    x = pfx + (box_w + gap) * i
    add_rect(s, x, pfy, box_w, Inches(1.2), WHITE,
             line=clr, rounded=True, radius=0.12)
    add_rect(s, x, pfy, box_w, Inches(0.3), clr, rounded=True, radius=0.5)
    add_text(s, x, pfy + Inches(0.02), box_w, Inches(0.3), str(i+1),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, pfy + Inches(0.4), box_w, Inches(0.45), ic,
             size=22, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, pfy + Inches(0.85), box_w, Inches(0.3), lab,
             size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    if i < len(pipe) - 1:
        arrow_x = x + box_w - Inches(0.02)
        add_text(s, arrow_x, pfy + Inches(0.45), Inches(0.2), Inches(0.3),
                 "›", size=22, bold=True, color=GRAY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 3 quality boosters (right side)
section_title(s, Inches(7.1), Inches(2.85), Inches(6), "品質向上の3系統", icon="🎚️")
triad = [
    ("🗣️", "VTT 話者情報", "Zoom 自動文字起こしの話者名を統合 → 担当者推定精度UP", TEAL),
    ("🖼️", "スライドOCR",   "固有名詞・技術用語・数値の誤変換を抑制",         CORAL),
    ("🤖", "Whisper + LLM", "高品質日本語ASR + 構造化抽出 (常時有効)",        SUCCESS),
]
ty = Inches(3.3)
for i, (ic, t, d, clr) in enumerate(triad):
    y = ty + Inches(1.15) * i
    add_rect(s, Inches(7.1), y, Inches(5.9), Inches(1.05), WHITE,
             line=RGBColor(0xD6, 0xE0, 0xEA), rounded=True, radius=0.1)
    add_rect(s, Inches(7.1), y, Inches(0.15), Inches(1.05), clr, rounded=True, radius=0.5)
    add_text(s, Inches(7.35), y + Inches(0.22), Inches(0.7), Inches(0.7), ic,
             size=22, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.1), y + Inches(0.18), Inches(4.8), Inches(0.38),
             t, size=13, bold=True, color=DARK)
    add_text(s, Inches(8.1), y + Inches(0.55), Inches(4.8), Inches(0.5),
             d, size=10, color=GRAY)

# Notes
section_title(s, Inches(7.1), Inches(6.8), Inches(6), "運用", icon="⚠️")
add_text(s, Inches(7.2), Inches(7.15), Inches(5.8), Inches(0.3),
         "進捗: スレッド公開  /  完了・エラー: ephemeral  /  同時1ジョブのみ  /  10–20分",
         size=9, color=GRAY)

footer(s, page=8, total=8)


prs.save(OUT)
print(f"Saved: {OUT}")
