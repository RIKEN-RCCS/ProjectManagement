#!/usr/bin/env python3
"""PPTX / PDF からスライド要約読み上げ付き mp4 を生成する。

各スライドについて:
  1. プレーンテキスト抽出 (PPTX→python-pptx, PDF→pdftotext / PyMuPDF)
  2. スライド画像生成 (LibreOffice→pdftoppm / PyMuPDF)
  3. 画像のマルチモーダル OCR (recording/slide_ocr.ocr_slide_image)
  4. (1)+(3) を併記して LLM で 2-3 文の narration に要約
  5. VOICEVOX (pm_tts.synth_chunk) で WAV 化
  6. ffmpeg で 静止画 + 音声 → セグメント mp4
全セグメントを ffmpeg concat demuxer で 1 本に結合する。

OCR は誤認識を含み得るので「抽出テキスト優先」を要約プロンプトで明示する。

呼び出し元は scripts/argus/pm_argus.py の _run_narrate(/argus-narrate)。
CLI は要約品質確認 (--dry-run) と動作確認用。
"""

from __future__ import annotations

import argparse
import logging
import os
import re
import shutil
import subprocess
import sys
import tempfile
import wave
from dataclasses import dataclass
from pathlib import Path

_SCRIPT_DIR = Path(__file__).resolve().parent
if str(_SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPT_DIR))

import pm_tts  # noqa: E402

logger = logging.getLogger(__name__)


# --------------------------------------------------------------------------- #
# 抽出
# --------------------------------------------------------------------------- #

@dataclass
class Slide:
    index: int          # 1-origin
    image_png: Path
    text: str           # 抽出テキスト (本文+notes)


def _extract_pptx_text(pptx_path: Path) -> list[str]:
    """PPTX の各スライドから本文+speaker notes を抽出する。"""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    out: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    parts.append(txt)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[ノート]\n{notes}")
        out.append("\n".join(parts))
    return out


def _extract_pdf_text(pdf_path: Path) -> list[str]:
    """PDF をページごとのテキスト list として返す。pdftotext → PyMuPDF の順で試す。"""
    try:
        text = subprocess.check_output(
            ["pdftotext", "-layout", "-enc", "UTF-8", str(pdf_path), "-"],
            timeout=120, text=True,
        )
        # pdftotext は \f (form feed) でページを区切る
        if "\f" in text:
            return [p.strip() for p in text.split("\f")]
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError):
        pass

    try:
        import fitz  # PyMuPDF
    except ImportError:
        return []
    try:
        doc = fitz.open(str(pdf_path))
        pages = [p.get_text().strip() for p in doc]
        doc.close()
        return pages
    except Exception as exc:
        logger.warning(f"PyMuPDF テキスト抽出失敗: {exc}")
        return []


def _prepare_pdf_and_images(src: Path, work_dir: Path) -> tuple[Path, list[Path]]:
    """src を PDF にし、ページ PNG リストと PDF パスを返す。pm_box_crawl の関数を再利用。"""
    from pm_box_crawl import _to_pdf, _pdf_to_images

    if src.suffix.lower() == ".pdf":
        pdf_path = src
    else:
        pdf_path = _to_pdf(src, work_dir)
        if pdf_path is None:
            raise RuntimeError(f"PDF 変換に失敗: {src}")

    images = _pdf_to_images(pdf_path, work_dir)
    if not images:
        raise RuntimeError(f"PDF 画像化に失敗: {pdf_path}")
    return pdf_path, images


def collect_slides(src: Path, work_dir: Path) -> list[Slide]:
    """PPTX / PDF からスライドのリストを構築する。画像とテキストの index は揃える。"""
    pdf_path, images = _prepare_pdf_and_images(src, work_dir)

    if src.suffix.lower() == ".pptx":
        texts = _extract_pptx_text(src)
    else:
        texts = _extract_pdf_text(pdf_path)

    # 枚数のずれは画像優先で揃える（OCR で補完できるため）
    n = len(images)
    if len(texts) < n:
        texts = texts + [""] * (n - len(texts))
    elif len(texts) > n:
        texts = texts[:n]

    return [
        Slide(index=i + 1, image_png=img, text=texts[i].strip())
        for i, img in enumerate(images)
    ]


# --------------------------------------------------------------------------- #
# 要約 (LLM)
# --------------------------------------------------------------------------- #

_NARRATION_SYSTEM_JA = (
    "あなたはスライド資料を読み上げ用に要約するアシスタントです。"
    "聴き手が短時間で資料の概要を把握できるよう、自然な日本語で簡潔にまとめてください。"
)

_NARRATION_SYSTEM_EN = (
    "You are an assistant that summarizes slide content for text-to-speech narration. "
    "Summarize each slide in concise, natural English so listeners can quickly grasp the key points."
)

_NARRATION_PROMPT_TMPL_JA = """\
このスライド({index}枚目)の内容を、音声読み上げ用に{max_sentences}文以内・合計{max_chars}文字以内の日本語平文に要約してください。

以下の 2 種類の入力があります。
- (A) スライドファイルから抽出したテキスト: 正確だが断片的なことがあります
- (B) 画像 OCR から得たテキスト: 誤認識を含む可能性があります

ルール:
- (A) を優先し、(A) で読み取れる固有名詞・数値・日付は維持する
- (A) が空または極めて短い場合のみ (B) を主体にする
- 箇条書き記号・URL・括弧書き・記号読みは出さない
- 「このスライドでは」のような前置きは不要、内容そのものから始める
- 出力は要約文のみ。マークダウン記号や見出しは付けない
- 内容が空・装飾のみの場合は「{index}枚目」とだけ書く

# (A) 抽出テキスト
{extracted}

# (B) OCR テキスト
{ocr}

要約:"""

_NARRATION_PROMPT_TMPL_EN = """\
Summarize slide {index} for text-to-speech narration in {max_sentences} sentences or fewer, \
within {max_chars} characters total.

Two input sources are provided:
- (A) Text extracted from the slide file: accurate but may be fragmented
- (B) Text from image OCR: may contain recognition errors

Rules:
- Prioritize (A); preserve proper nouns, numbers, and dates exactly as they appear in (A)
- Use (B) as the primary source only if (A) is empty or very short
- Omit bullet markers, URLs, parenthetical text, and symbol readings
- Start directly with the content — no preamble like "This slide shows..."
- Output the narration text only. No markdown or headings.
- If the slide has no meaningful content, write only "Slide {index}."

# (A) Extracted text
{extracted}

# (B) OCR text
{ocr}

Narration:"""


def _summarize_slide(
    slide: Slide,
    ocr_text: str,
    *,
    max_sentences: int,
    max_chars: int,
    lang: str = "ja",
    timeout: int = 60,
) -> str:
    from cli_utils import call_argus_llm, strip_think_blocks

    if lang == "en":
        system = _NARRATION_SYSTEM_EN
        tmpl = _NARRATION_PROMPT_TMPL_EN
        empty_marker = "(none)"
        fallback_label = f"Slide {slide.index}."
    else:
        system = _NARRATION_SYSTEM_JA
        tmpl = _NARRATION_PROMPT_TMPL_JA
        empty_marker = "(なし)"
        fallback_label = f"{slide.index}枚目"

    extracted = slide.text or empty_marker
    ocr = (ocr_text or "").strip() or empty_marker

    prompt = tmpl.format(
        index=slide.index,
        extracted=extracted[:4000],
        ocr=ocr[:4000],
        max_sentences=max_sentences,
        max_chars=max_chars,
    )

    try:
        raw = call_argus_llm(
            prompt,
            system=system,
            max_tokens=512,
            timeout=timeout,
        )
        out = strip_think_blocks(raw).strip()
    except Exception as exc:
        logger.warning(f"slide {slide.index}: LLM 要約失敗 ({exc}) — 抽出テキスト先頭を使用")
        fallback = (slide.text or ocr_text or "").strip().splitlines()
        head = " ".join(s for s in fallback if s)[:max_chars]
        return head or fallback_label

    out = re.sub(r"^(要約[::]\s*|出力[::]\s*|Narration[::]\s*)", "", out)
    if not out:
        return fallback_label
    if not out.endswith(("。", "！", "？", ".", "!", "?")):
        out += "." if lang == "en" else "。"
    return out


# --------------------------------------------------------------------------- #
# OCR
# --------------------------------------------------------------------------- #

def _ocr_slide(image: Path) -> str:
    base_url = os.environ.get("OPENAI_API_BASE")
    if not base_url:
        return ""
    try:
        from recording.slide_ocr import MEETING_FRAME_OCR_PROMPT
        from pm_box_crawl import ocr_slide_image
    except ImportError as exc:
        logger.warning(f"OCR モジュールの import に失敗: {exc}")
        return ""

    try:
        md = ocr_slide_image(image, base_url, prompt=MEETING_FRAME_OCR_PROMPT)
    except Exception as exc:
        logger.warning(f"OCR 失敗 ({image.name}): {exc}")
        return ""
    if not md:
        return ""
    stripped = md.strip()
    head = stripped.splitlines()[0].strip() if stripped else ""
    if any(marker in head for marker in ("[テキストなし]", "[会議無関係]")):
        return ""
    return stripped


# --------------------------------------------------------------------------- #
# TTS / 動画組み立て
# --------------------------------------------------------------------------- #

def _ffmpeg_or_raise() -> str:
    path = shutil.which("ffmpeg")
    if not path:
        raise RuntimeError("ffmpeg が PATH に見つかりません")
    return path


def _trim_silence(in_wav: Path, out_wav: Path, pad_sec: float = 0.15) -> None:
    """先頭・末尾の無音を ffmpeg silenceremove で除去し、両端に短いパディングを付ける。

    スライド境界の不自然な空白（VOICEVOX の pre/postPhonemeLength と文末記号の余韻が
    スライド N の末尾 + N+1 の先頭で重なって 5 秒近くになる現象）を抑える。
    """
    _ffmpeg_or_raise()
    af = (
        # 先頭の無音を除去
        "silenceremove=start_periods=1:start_duration=0:start_threshold=-45dB,"
        # areverse で末尾を先頭にし、再度 silenceremove で末尾無音を除去
        "areverse,"
        "silenceremove=start_periods=1:start_duration=0:start_threshold=-45dB,"
        "areverse,"
        # 両端に短いパディングを足す
        f"adelay={int(pad_sec * 1000)}|{int(pad_sec * 1000)},"
        f"apad=pad_dur={pad_sec}"
    )
    cmd = [
        "ffmpeg", "-y", "-loglevel", "error",
        "-i", str(in_wav),
        "-af", af,
        str(out_wav),
    ]
    subprocess.run(cmd, check=True)


def _synth_narration_to_wav(
    text: str,
    out_wav: Path,
    *,
    speaker: int,
    speed: float,
    chunk_limit: int = pm_tts.VOICEVOX_TEXT_LIMIT,
) -> None:
    """1 スライド分の narration を WAV にする (チャンク合成 → concat → 無音トリム)。"""
    chunks = pm_tts.split_into_sentences(text, limit=chunk_limit)
    if not chunks:
        chunks = [text]
    tmp = Path(tempfile.mkdtemp(prefix="narr_chunks_"))
    try:
        wavs: list[Path] = []
        for i, c in enumerate(chunks, 1):
            wp = tmp / f"chunk_{i:04d}.wav"
            pm_tts.synth_chunk(c, speaker, wp, speed=speed)
            wavs.append(wp)
        merged = tmp / "merged.wav"
        pm_tts.concat_wavs(wavs, merged)
        _trim_silence(merged, out_wav)
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def _wav_duration(path: Path) -> float:
    with wave.open(str(path), "rb") as w:
        frames = w.getnframes()
        rate = w.getframerate()
    return frames / float(rate) if rate else 0.0


def _mux_video(
    image_wav_pairs: list[tuple[Path, Path]],
    out_mp4: Path,
    work_dir: Path,
) -> None:
    """画像列 + 音声列 → 1 本の mp4。

    各画像の表示時間を対応 WAV の正確な秒数にロックすることで、
    セグメント mp4 → concat 方式で発生していた「fps 量子化による
    画像と音声の累積ずれ」を排除する。

    - 画像は ffmpeg concat demuxer の `file` + `duration` で連続化
    - 音声は wav を concat demuxer で連結
    - 最後に video / audio を 1 回だけ mux （faststart 付与）
    """
    _ffmpeg_or_raise()

    img_list = work_dir / "images.txt"
    aud_list = work_dir / "audio.txt"

    img_lines: list[str] = []
    aud_lines: list[str] = []
    for image, wav in image_wav_pairs:
        dur = _wav_duration(wav)
        img_lines.append(f"file '{image.resolve()}'")
        img_lines.append(f"duration {dur:.6f}")
        aud_lines.append(f"file '{wav.resolve()}'")
    # concat demuxer 仕様: 末尾フレームを保持するため最後の file を duration なしで再記述
    img_lines.append(f"file '{image_wav_pairs[-1][0].resolve()}'")

    img_list.write_text("\n".join(img_lines) + "\n", encoding="utf-8")
    aud_list.write_text("\n".join(aud_lines) + "\n", encoding="utf-8")

    # 静止画用に 5 fps CFR で出力する。concat demuxer の duration は CFR 化で
    # 最寄りフレームに snap されるが、各画像は独立に snap されるだけで誤差は
    # 累積しない (±0.1 秒以内)。
    # キーフレームはスライド境界に 1 個入れば十分で、x264 のデフォルト scenecut
    # detection (threshold=40) がスライド切り替えを自動検知して I フレームを
    # 挿入する。`-g` は保険の最大間隔として 60 秒 (5fps × 60) に設定する。
    # `-keyint_min` を低くすると scenecut が抑制されるので 1 のままにする。
    cmd = [
        "ffmpeg", "-y", "-loglevel", "error",
        "-f", "concat", "-safe", "0", "-i", str(img_list),
        "-f", "concat", "-safe", "0", "-i", str(aud_list),
        "-map", "0:v", "-map", "1:a",
        "-fps_mode", "cfr",
        "-c:v", "libx264", "-tune", "stillimage", "-pix_fmt", "yuv420p",
        "-vf", "scale=trunc(iw/2)*2:trunc(ih/2)*2",
        "-r", "5", "-g", "300", "-keyint_min", "1",
        "-c:a", "aac", "-b:a", "128k", "-ar", "44100",
        "-movflags", "+faststart",
        "-shortest",
        str(out_mp4),
    ]
    subprocess.run(cmd, check=True)


# --------------------------------------------------------------------------- #
# 高レベル API
# --------------------------------------------------------------------------- #

def build_slide_video(
    src: Path,
    output_mp4: Path,
    *,
    speaker: int = pm_tts.DEFAULT_SPEAKER,
    speed: float = 1.0,
    max_sentences: int = 3,
    max_chars: int = 180,
    max_slides: int | None = None,
    lang: str = "ja",
    quiet: bool = False,
) -> Path:
    """PPTX/PDF → mp4 を生成して output_mp4 のパスを返す。"""
    src = src.resolve()
    if not src.is_file():
        raise FileNotFoundError(src)
    suffix = src.suffix.lower()
    if suffix not in (".pptx", ".pdf"):
        raise ValueError(f"対応形式は .pptx / .pdf のみ: {src.name}")

    output_mp4 = output_mp4.resolve()
    output_mp4.parent.mkdir(parents=True, exist_ok=True)

    pm_tts.check_tts_alive()
    _ffmpeg_or_raise()

    work_root = Path(tempfile.mkdtemp(prefix="slide_video_"))
    try:
        slides = collect_slides(src, work_root)
        if max_slides is not None and len(slides) > max_slides:
            slides = slides[:max_slides]
        if not slides:
            raise RuntimeError("スライドが 0 枚です")

        logger.info(f"スライド {len(slides)} 枚を処理")

        seg_dir = work_root / "segments"
        seg_dir.mkdir(parents=True, exist_ok=True)

        pairs: list[tuple[Path, Path]] = []
        for slide in pm_tts._iter_progress(slides, total=len(slides), desc="slides", quiet=quiet):
            ocr_text = _ocr_slide(slide.image_png)
            narration = _summarize_slide(
                slide, ocr_text,
                max_sentences=max_sentences,
                max_chars=max_chars,
                lang=lang,
            )
            logger.info(f"  slide {slide.index}: narration={narration}")

            wav_path = seg_dir / f"slide_{slide.index:04d}.wav"
            _synth_narration_to_wav(narration, wav_path, speaker=speaker, speed=speed)
            pairs.append((slide.image_png, wav_path))

        _mux_video(pairs, output_mp4, work_root)
    finally:
        shutil.rmtree(work_root, ignore_errors=True)

    return output_mp4


# --------------------------------------------------------------------------- #
# CLI
# --------------------------------------------------------------------------- #

def main() -> int:
    ap = argparse.ArgumentParser(description="PPTX/PDF を要約読み上げ付き mp4 に変換する")
    ap.add_argument("input", type=Path, help="入力 .pptx / .pdf")
    ap.add_argument("-o", "--output", type=Path, default=None,
                    help="出力 mp4（既定: 入力と同じ basename .mp4）")
    ap.add_argument("--speaker", type=int, default=pm_tts.DEFAULT_SPEAKER)
    ap.add_argument("--speed", type=float, default=1.0)
    ap.add_argument("--max-sentences", type=int, default=3)
    ap.add_argument("--max-chars", type=int, default=180)
    ap.add_argument("--max-slides", type=int, default=None,
                    help="先頭 N 枚のみ処理 (動作確認用)")
    ap.add_argument("--lang", choices=["ja", "en"], default="ja",
                    help="ナレーション言語 (ja=日本語, en=英語, デフォルト: ja)")
    ap.add_argument("--dry-run", action="store_true",
                    help="narration 文だけ標準出力に流して mp4 は作らない")
    ap.add_argument("--verbose", "-v", action="store_true")
    args = ap.parse_args()

    logging.basicConfig(
        level=logging.INFO if args.verbose else logging.WARNING,
        format="%(asctime)s %(levelname)s %(message)s",
        stream=sys.stderr,
    )

    if not args.input.is_file():
        print(f"入力が存在しません: {args.input}", file=sys.stderr)
        return 1

    if args.dry_run:
        work = Path(tempfile.mkdtemp(prefix="slide_video_dry_"))
        try:
            slides = collect_slides(args.input.resolve(), work)
            if args.max_slides:
                slides = slides[: args.max_slides]
            for slide in slides:
                ocr_text = _ocr_slide(slide.image_png)
                narration = _summarize_slide(
                    slide, ocr_text,
                    max_sentences=args.max_sentences,
                    max_chars=args.max_chars,
                    lang=args.lang,
                )
                print(f"--- slide {slide.index} ---")
                print(narration)
        finally:
            shutil.rmtree(work, ignore_errors=True)
        return 0

    output_mp4 = args.output or args.input.with_suffix(".mp4")
    build_slide_video(
        args.input,
        output_mp4,
        speaker=args.speaker,
        speed=args.speed,
        max_sentences=args.max_sentences,
        max_chars=args.max_chars,
        max_slides=args.max_slides,
        lang=args.lang,
    )
    size_mb = output_mp4.stat().st_size / (1024 * 1024)
    print(f"wrote {output_mp4} ({size_mb:.2f} MB)", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
