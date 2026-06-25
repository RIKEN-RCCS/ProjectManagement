"""Ocean Gradient テーマの共通描画ヘルパー。

pptx 生成スクリプトから再利用できる色・フォント・基本シェイプ関数を提供する。
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---- Ocean Gradient palette ----
NAVY    = RGBColor(0x0B, 0x2A, 0x4A)
DEEP    = RGBColor(0x06, 0x5A, 0x82)
TEAL    = RGBColor(0x1C, 0x72, 0x93)
MINT    = RGBColor(0x9A, 0xD1, 0xD4)
ICE     = RGBColor(0xEE, 0xF5, 0xF9)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x16, 0x1F, 0x2C)
GRAY    = RGBColor(0x5A, 0x6A, 0x7E)
MUTED   = RGBColor(0x9A, 0xA7, 0xB8)
CORAL   = RGBColor(0xF1, 0x6B, 0x4F)
GOLD    = RGBColor(0xE8, 0xB5, 0x3B)
SUCCESS = RGBColor(0x4C, 0xAF, 0x50)
DANGER  = RGBColor(0xD1, 0x4B, 0x4B)
LINE    = RGBColor(0xD6, 0xE0, 0xEA)

HEADER_FONT = "Meiryo"
BODY_FONT   = "Meiryo"
MONO_FONT   = "Consolas"


def make_presentation(width_in: float = 13.333, height_in: float = 7.5) -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(width_in)
    prs.slide_height = Inches(height_in)
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---- primitives ----
def add_rect(slide, x, y, w, h, color, *, line=None, rounded=False, radius=0.12):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line; shape.line.width = Pt(0.75)
    if rounded:
        shape.adjustments[0] = radius
    return shape


def add_bg(slide, sw, sh, color=WHITE):
    return add_rect(slide, 0, 0, sw, sh, color)


def add_text(slide, x, y, w, h, text, *, size=16, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font=BODY_FONT, anchor=None, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_paragraph_text(slide, x, y, w, h, text, *, size=12, color=DARK,
                       font=BODY_FONT, line_spacing=1.25,
                       paragraph_gap_pt=6, heading_color=None,
                       heading_pattern=r"【([^】]+)】"):
    """長い段落テキストを読みやすくレイアウトする。

    - 改行ごとに段落分け
    - paragraph_gap_pt で段落間にスペース
    - 行内に【見出し】があれば heading_color で太字化
    - line_spacing は単位なしの相対倍率（pptx の Paragraph.line_spacing 仕様）
    """
    import re as _re
    heading_color = heading_color or NAVY
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    pat = _re.compile(heading_pattern)

    paragraphs = [p for p in (text or "").split("\n") if p.strip()]
    for i, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(paragraph_gap_pt)
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass
        m = pat.match(line.strip())
        if m:
            head = m.group(0)
            rest = line.strip()[len(head):]
            r1 = p.add_run(); r1.text = head
            r1.font.name = font; r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = heading_color
            if rest:
                r2 = p.add_run(); r2.text = rest
                r2.font.name = font; r2.font.size = Pt(size); r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = line
            r.font.name = font; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=DARK,
                bullet_color=None, gap=4, marker="▸ "):
    bullet_color = bullet_color or TEAL
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_top = Inches(0.02)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        r1 = p.add_run(); r1.text = marker
        r1.font.name = BODY_FONT; r1.font.size = Pt(size); r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run(); r2.text = it
        r2.font.name = BODY_FONT; r2.font.size = Pt(size); r2.font.color.rgb = color
    return tb


def code_block(slide, x, y, w, h, lines, *, size=12):
    add_rect(slide, x, y, w, h, RGBColor(0x10, 0x20, 0x30), rounded=True, radius=0.05)
    for i, c in enumerate([CORAL, GOLD, SUCCESS]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.18 + i*0.22), y + Inches(0.15),
                                     Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = c
        dot.line.fill.background()
    tb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.4),
                                  w - Inches(0.5), h - Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run(); r.text = ln
        r.font.name = MONO_FONT; r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0xDA, 0xE8, 0xF4)


def number_badge(slide, x, y, num, *, size_in=0.7, bg=TEAL, fg=WHITE):
    d = Inches(size_in)
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    c.fill.solid(); c.fill.fore_color.rgb = bg
    c.line.fill.background()
    add_text(slide, x, y, d, d, str(num), size=int(size_in * 28),
             bold=True, color=fg, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font=HEADER_FONT)


def header(slide, sh, number, title, subtitle=None):
    """共通ヘッダ（左 navy バー + 番号バッジ + タイトル + 区切り線）。

    sh は slide_height（Inches で構築済みの値）。
    """
    add_rect(slide, 0, 0, Inches(0.35), sh, NAVY)
    add_rect(slide, Inches(0.35), 0, Inches(0.06), sh, TEAL)
    add_rect(slide, Inches(0.55), Inches(0.45), Inches(12.5), Inches(0.05), MINT)
    number_badge(slide, Inches(0.65), Inches(0.65), number, size_in=0.85)
    add_text(slide, Inches(1.75), Inches(0.65), Inches(11.2), Inches(0.55),
             title, size=24, bold=True, color=NAVY, font=HEADER_FONT)
    if subtitle:
        add_text(slide, Inches(1.75), Inches(1.20), Inches(11.2), Inches(0.4),
                 subtitle, size=13, color=GRAY, italic=True)
    add_rect(slide, Inches(0.55), Inches(1.72), Inches(12.5), Inches(0.02),
             RGBColor(0xE0, 0xE8, 0xF0))


def section_title(slide, x, y, w, text):
    add_rect(slide, x, y + Inches(0.05), Inches(0.08), Inches(0.32), CORAL)
    x2 = x + Inches(0.2)
    w2 = w - Inches(0.2)
    add_text(slide, x2, y, w2, Inches(0.4), text, size=14, bold=True,
             color=TEAL, font=HEADER_FONT)


def callout(slide, x, y, w, h, label, text, *,
            bg=RGBColor(0xFF, 0xF6, 0xE5), border=GOLD,
            label_color=RGBColor(0x8A, 0x5A, 0x00), text_size=12):
    add_rect(slide, x, y, w, h, bg, line=border, rounded=True, radius=0.06)
    add_text(slide, x + Inches(0.2), y + Inches(0.12), w - Inches(0.4), Inches(0.35),
             label, size=12, bold=True, color=label_color)
    add_text(slide, x + Inches(0.2), y + Inches(0.48), w - Inches(0.4), h - Inches(0.55),
             text, size=text_size, color=DARK)


def footer(slide, sw, sh, text="", page=None, total=None):
    add_text(slide, Inches(0.55), sh - Inches(0.35), Inches(8), Inches(0.25),
             text, size=9, color=MUTED)
    if page is not None:
        add_text(slide, sw - Inches(1.2), sh - Inches(0.35), Inches(0.8), Inches(0.25),
                 f"{page} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def progress_bar(slide, x, y, w, h, ratio, *, fill_color=None, bg_color=None):
    """0.0〜1.0 の達成率バー。fill_color は ratio に応じて自動選択も可。"""
    bg_color = bg_color or RGBColor(0xE6, 0xEC, 0xF2)
    if fill_color is None:
        if ratio >= 0.7:
            fill_color = SUCCESS
        elif ratio >= 0.4:
            fill_color = GOLD
        else:
            fill_color = CORAL
    add_rect(slide, x, y, w, h, bg_color, rounded=True, radius=0.5)
    fw = max(int(w * max(0.0, min(1.0, ratio))), 1)
    if fw > 0:
        add_rect(slide, x, y, fw, h, fill_color, rounded=True, radius=0.5)
