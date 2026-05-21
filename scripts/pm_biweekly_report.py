#!/usr/bin/env python3
"""pm_biweekly_report.py — 隔週進捗会議向け pptx レポート生成。

マイルストーン（goals.yaml 由来）を縦軸に、期間内の動き（完了/新規 AI、決定事項、
確定ナレッジ、Slack 議論、BOX 資料）をまとめ、LLM で短いナラティブを添えて pptx
スライドに出す。

使い方:
    python3 scripts/pm_biweekly_report.py --since 2026-05-06 --until 2026-05-19
    python3 scripts/pm_biweekly_report.py --index-name pm-hpc --output reports/biweekly.pptx
    python3 scripts/pm_biweekly_report.py --markdown-only --dry-run
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from collections import defaultdict
from dataclasses import dataclass, field
from datetime import date, datetime, timedelta
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO_ROOT / "scripts"))

import yaml  # type: ignore

from db_utils import (
    open_pm_db,
    open_db,
    open_knowledge_db,
    fetch_milestone_progress,
)
from cli_utils import (call_argus_llm, strip_think_blocks,  # noqa: E402
                       add_filter_arg, resolve_filter_presets)

from pptx.util import Inches, Pt  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402

import pptx_theme as T  # noqa: E402


# --------------------------------------------------------------------------- #
# データクラス
# --------------------------------------------------------------------------- #
@dataclass
class Milestone:
    milestone_id: str
    goal_id: str
    name: str
    due_date: str | None
    area: str | None
    success_criteria: list[str]
    open_count: int = 0
    closed_count: int = 0

    keywords: list[str] = field(default_factory=list)
    completed_in_period: list[dict] = field(default_factory=list)
    created_in_period: list[dict] = field(default_factory=list)
    decisions_in_period: list[dict] = field(default_factory=list)
    knowledge_in_period: list[dict] = field(default_factory=list)
    box_in_period: list[dict] = field(default_factory=list)
    slack_threads_in_period: list[dict] = field(default_factory=list)
    upcoming_due: list[dict] = field(default_factory=list)
    narrative: str = ""

    def progress_ratio(self) -> float:
        total = self.open_count + self.closed_count
        return (self.closed_count / total) if total else 0.0

    def days_remaining(self) -> int | None:
        if not self.due_date:
            return None
        try:
            d = datetime.strptime(self.due_date, "%Y-%m-%d").date()
            return (d - date.today()).days
        except ValueError:
            return None


@dataclass
class RiskSignals:
    """機械的に検出した障害シグナル（LLM へ「観察事実」として渡す）。"""
    overdue_total: int = 0
    overdue_by_category: dict = field(default_factory=dict)  # {"年度計画策定": [...], "外部依存": [...], "その他": [...]}
    overdue_oldest_days: int = 0
    overdue_undefined_assignee: int = 0  # assignee が None / "未定" / "null" / 空のもの

    no_milestone_link_ratio: float = 0.0  # 期間内 AI のうち milestone 未紐付け率
    no_milestone_count: int = 0

    progress_zero_milestones: list[str] = field(default_factory=list)  # 紐付き AI 0/0 の milestone ID
    short_deadline_milestones: list[str] = field(default_factory=list)  # 残30日以内なのに紐付き 0 件のもの

    assignee_concentration: list[tuple[str, int]] = field(default_factory=list)  # [(name, count)]
    undefined_assignee_in_period: int = 0  # 期間内 AI で未定担当の数

    external_dep_signals: list[str] = field(default_factory=list)  # 外部依存（NVIDIA / 富士通 / 契約 等）の文脈で検出された AI/decision の抜粋
    fy_planning_overdue: list[str] = field(default_factory=list)  # FY26 計画策定系の期限超過 AI の抜粋

    def has_critical(self) -> bool:
        return (self.fy_planning_overdue
                or self.short_deadline_milestones
                or self.overdue_total >= 20)


@dataclass
class ReportData:
    since: str
    until: str
    index_name: str
    milestones: list[Milestone]
    unlinked_completed: list[dict]
    unlinked_created: list[dict]
    unlinked_decisions: list[dict]
    overdue_items: list[dict]
    summary_narrative: str = ""
    risks: RiskSignals = field(default_factory=RiskSignals)


# --------------------------------------------------------------------------- #
# キーワード抽出（milestone 名 + success_criteria + area からトークンを作る）
# --------------------------------------------------------------------------- #
_NOISE = {
    "の", "を", "に", "は", "が", "で", "と", "や", "へ", "から", "まで", "より",
    "ある", "する", "なる", "おる", "いる", "こと", "もの", "ため", "ここ", "そこ",
    "これ", "それ", "あれ", "どれ", "場合", "等", "など", "及び", "または",
    "について", "に対して", "に関する", "における", "として",
    "the", "a", "an", "and", "or", "of", "to", "in", "on", "for", "by", "with",
    "について", "完了", "確立", "受領", "提出", "整理", "策定", "決定", "実施",
    "対応", "確定", "明文化", "運用", "開催", "公開", "確認", "報告",
}


def extract_keywords(milestone: Milestone) -> list[str]:
    text = " ".join(filter(None, [
        milestone.name,
        " ".join(milestone.success_criteria or []),
        milestone.area or "",
    ]))
    # 英数字・カタカナ・漢字の連続を抽出
    tokens = re.findall(r"[A-Za-z0-9][A-Za-z0-9\-_/\.+]{2,}|[゠-ヿ]{3,}|[一-鿿]{2,}", text)
    seen: set[str] = set()
    keywords: list[str] = []
    for t in tokens:
        t_lower = t.lower()
        if t_lower in _NOISE or t in _NOISE:
            continue
        if t in seen:
            continue
        seen.add(t)
        keywords.append(t)
    return keywords


def keyword_match(text: str, keywords: list[str]) -> int:
    if not text:
        return 0
    score = 0
    for kw in keywords:
        if kw.lower() in text.lower():
            score += 1
    return score


# --------------------------------------------------------------------------- #
# ID パーサー
# --------------------------------------------------------------------------- #
def assign_to_milestone(item_text: str, milestones: list[Milestone],
                       explicit_id: str | None = None) -> Milestone | None:
    """item を milestone に紐づける。explicit_id 優先。なければキーワードマッチ。"""
    if explicit_id:
        for m in milestones:
            if m.milestone_id == explicit_id:
                return m
    best, best_score = None, 0
    for m in milestones:
        s = keyword_match(item_text, m.keywords)
        if s > best_score:
            best, best_score = m, s
    # 閾値: 2語以上一致したものだけ自動紐付け
    return best if best_score >= 2 else None


# --------------------------------------------------------------------------- #
# 期間内データの取得
# --------------------------------------------------------------------------- #
def _build_source_filter(channel_ids: list[str] | None,
                          meeting_kinds: list[str] | None,
                          ai_alias: str = "a",
                          meetings_join: bool = False) -> tuple[str, list]:
    """channel_ids / meeting_kinds から WHERE 句の部分文字列とパラメータを返す。

    meetings_join=True の場合は既に meetings テーブルが m として JOIN されている前提。
    False の場合は channel_id のみのフィルタ（meetings JOIN がないクエリ用）。
    """
    if not channel_ids and not meeting_kinds:
        return "", []
    clauses: list[str] = []
    params: list = []
    if channel_ids:
        placeholders = ",".join("?" * len(channel_ids))
        clauses.append(f"{ai_alias}.channel_id IN ({placeholders})")
        params.extend(channel_ids)
    if meeting_kinds and meetings_join:
        placeholders = ",".join("?" * len(meeting_kinds))
        clauses.append(f"m.kind IN ({placeholders})")
        params.extend(meeting_kinds)
    elif meeting_kinds and not meetings_join:
        placeholders = ",".join("?" * len(meeting_kinds))
        clauses.append(
            f"{ai_alias}.meeting_id IN "
            f"(SELECT meeting_id FROM meetings WHERE kind IN ({placeholders}))"
        )
        params.extend(meeting_kinds)
    return " AND (" + " OR ".join(clauses) + ")", params


def fetch_period_action_items(conn, since: str, until: str,
                              channel_ids: list[str] | None = None,
                              meeting_kinds: list[str] | None = None,
                              ) -> tuple[list[dict], list[dict]]:
    """期間内に完了した AI / 期間内に新規発生した AI を取得"""
    src_filter, src_params = _build_source_filter(channel_ids, meeting_kinds,
                                                   ai_alias="a", meetings_join=False)
    completed = conn.execute(
        f"""
        SELECT a.id, a.content, a.assignee, a.due_date, a.milestone_id,
               a.source, a.source_ref, a.extracted_at,
               (SELECT MAX(changed_at) FROM audit_log
                WHERE table_name='action_items' AND record_id=CAST(a.id AS TEXT)
                  AND field='status' AND new_value='closed') AS closed_at
        FROM action_items a
        WHERE a.status='closed' AND COALESCE(a.deleted,0)=0{src_filter}
        """,
        src_params,
    ).fetchall()
    completed_in = []
    for r in completed:
        d = dict(r)
        ts = d.get("closed_at") or d.get("extracted_at")
        if ts and since <= ts[:10] <= until:
            d["_period_date"] = ts[:10]
            completed_in.append(d)

    created = conn.execute(
        f"""
        SELECT id, content, assignee, due_date, milestone_id,
               source, source_ref, extracted_at, status, channel_id, meeting_id
        FROM action_items a
        WHERE COALESCE(deleted,0)=0 AND extracted_at >= ? AND extracted_at <= ?{src_filter}
        """,
        [since, until, *src_params],
    ).fetchall()
    created_in = [dict(r) for r in created]
    for d in created_in:
        d["_period_date"] = (d.get("extracted_at") or "")[:10]
    return completed_in, created_in


def fetch_period_decisions(conn, since: str, until: str,
                           channel_ids: list[str] | None = None,
                           meeting_kinds: list[str] | None = None) -> list[dict]:
    src_filter, src_params = _build_source_filter(channel_ids, meeting_kinds,
                                                   ai_alias="d", meetings_join=False)
    rows = conn.execute(
        f"""
        SELECT id, content, decided_at, source, source_ref, decided_by, rationale,
               channel_id, meeting_id
        FROM decisions d
        WHERE COALESCE(deleted,0)=0 AND decided_at >= ? AND decided_at <= ?{src_filter}
        ORDER BY decided_at DESC
        """,
        [since, until, *src_params],
    ).fetchall()
    return [dict(r) for r in rows]


def fetch_upcoming_due(conn, until: str, lookahead_days: int,
                       channel_ids: list[str] | None = None,
                       meeting_kinds: list[str] | None = None) -> list[dict]:
    horizon = (datetime.strptime(until, "%Y-%m-%d").date() + timedelta(days=lookahead_days)).isoformat()
    src_filter, src_params = _build_source_filter(channel_ids, meeting_kinds,
                                                   ai_alias="a", meetings_join=False)
    rows = conn.execute(
        f"""
        SELECT id, content, assignee, due_date, milestone_id, channel_id, meeting_id
        FROM action_items a
        WHERE status='open' AND COALESCE(deleted,0)=0
          AND due_date IS NOT NULL
          AND due_date > ? AND due_date <= ?{src_filter}
        ORDER BY due_date ASC
        """,
        [until, horizon, *src_params],
    ).fetchall()
    return [dict(r) for r in rows]


def fetch_overdue(conn, until: str,
                  channel_ids: list[str] | None = None,
                  meeting_kinds: list[str] | None = None) -> list[dict]:
    src_filter, src_params = _build_source_filter(channel_ids, meeting_kinds,
                                                   ai_alias="a", meetings_join=False)
    rows = conn.execute(
        f"""
        SELECT id, content, assignee, due_date, milestone_id, channel_id, meeting_id
        FROM action_items a
        WHERE status='open' AND COALESCE(deleted,0)=0
          AND due_date IS NOT NULL AND due_date < ?{src_filter}
        ORDER BY due_date ASC
        """,
        [until, *src_params],
    ).fetchall()
    return [dict(r) for r in rows]


def fetch_period_knowledge(knowledge_db: Path, since: str, until: str) -> list[dict]:
    if not knowledge_db.exists():
        return []
    try:
        k = open_knowledge_db(knowledge_db)
    except Exception:
        return []
    try:
        rows = k.execute(
            """
            SELECT id, kind, topic, current_state, rationale, tags,
                   confidence, decided_at, last_validated_at, created_at
            FROM knowledge
            WHERE COALESCE(deleted,0)=0 AND superseded_by IS NULL
              AND substr(created_at,1,10) >= ? AND substr(created_at,1,10) <= ?
            ORDER BY created_at DESC
            """,
            (since, until),
        ).fetchall()
        return [dict(r) for r in rows]
    finally:
        k.close()


def fetch_period_box_files(box_db: Path, since: str, until: str,
                           index_filter: list[str] | None = None) -> list[dict]:
    if not box_db.exists():
        return []
    try:
        b = open_db(box_db)
    except Exception:
        return []
    try:
        # box_docs.db の box_files テーブルを参照。modified_at は 'YYYY-MM-DD HH:MM' 形式。
        query = """
            SELECT box_file_id, name, file_format, modified_at, folder_path,
                   index_name, source_name, relevance
            FROM box_files
            WHERE substr(modified_at,1,10) >= ? AND substr(modified_at,1,10) <= ?
              AND COALESCE(relevance,'') != 'noise'
        """
        params: list = [since, until]
        if index_filter:
            placeholders = ",".join(["?"] * len(index_filter))
            query += f" AND index_name IN ({placeholders})"
            params.extend(index_filter)
        query += " ORDER BY modified_at DESC"
        rows = b.execute(query, params).fetchall()
        return [dict(r) for r in rows]
    finally:
        b.close()


def fetch_period_slack_threads(slack_db: Path, since: str, until: str,
                               channels: list[str], limit: int = 200) -> list[dict]:
    if not slack_db.exists() or not channels:
        return []
    try:
        s = open_db(slack_db)
    except Exception:
        return []
    try:
        placeholders = ",".join(["?"] * len(channels))
        # messages.timestamp は JST 'YYYY-MM-DD HH:MM:SS'
        rows = s.execute(
            f"""
            SELECT thread_ts, channel_id, user_name, text, timestamp, permalink
            FROM messages
            WHERE channel_id IN ({placeholders})
              AND substr(timestamp,1,10) >= ? AND substr(timestamp,1,10) <= ?
              AND text IS NOT NULL AND length(text) > 30
            ORDER BY timestamp DESC
            LIMIT ?
            """,
            [*channels, since, until, limit],
        ).fetchall()
        return [dict(r) for r in rows]
    finally:
        s.close()


# --------------------------------------------------------------------------- #
# argus_config.yaml
# --------------------------------------------------------------------------- #
def load_argus_config(path: Path) -> dict:
    with open(path, encoding="utf-8") as f:
        return yaml.safe_load(f) or {}


def resolve_index(cfg: dict, index_name: str) -> dict:
    indices = cfg.get("indices") or {}
    if index_name not in indices:
        raise SystemExit(f"index_name '{index_name}' は argus_config.yaml に定義されていません")
    entry = indices[index_name] or {}
    return {
        "channels": [c.strip() for c in (entry.get("channels") or []) if c and c.strip()],
        "minutes": entry.get("minutes") or [],
    }


# --------------------------------------------------------------------------- #
# リスクシグナル収集（LLM 入力用の機械的事実）
# --------------------------------------------------------------------------- #
_FY_PLANNING_PAT = re.compile(
    r"FY20\d\d|2026年度|今年度|年度計画|目標およびマイルストーン|目標.{0,5}マイルストーン|スケジュール.*更新|"
    r"見積.*リソース|計算リソース.*見積|目次案|計画報告書",
    re.IGNORECASE,
)
_EXTERNAL_DEP_PAT = re.compile(
    r"NVIDIA|富士通|Fujitsu|契約|MoU|NDA|機密|TSMC|BRCM|Synopsys|外部開発者|ベンダ|受領|"
    r"アップストリーム|上流",
    re.IGNORECASE,
)
_UNDEFINED_ASSIGNEES = {"", "未定", "null", "None", "-", "TBD", "tbd"}


def _is_undefined(assignee) -> bool:
    if assignee is None:
        return True
    s = str(assignee).strip()
    return s in _UNDEFINED_ASSIGNEES or s.lower() == "null"


def _categorize_overdue(content: str) -> str:
    if _FY_PLANNING_PAT.search(content or ""):
        return "年度計画策定"
    if _EXTERNAL_DEP_PAT.search(content or ""):
        return "外部依存"
    return "その他"


def collect_risk_signals(rd: "ReportData", today: str) -> RiskSignals:
    rs = RiskSignals()
    rs.overdue_total = len(rd.overdue_items)

    # 期限超過のカテゴリ分類 + 最古の超過日数
    today_dt = datetime.strptime(today, "%Y-%m-%d").date() if today else date.today()
    by_cat: dict[str, list[dict]] = {"年度計画策定": [], "外部依存": [], "その他": []}
    oldest_delta = 0
    undefined_overdue = 0
    for a in rd.overdue_items:
        cat = _categorize_overdue(a.get("content") or "")
        by_cat[cat].append(a)
        if _is_undefined(a.get("assignee")):
            undefined_overdue += 1
        try:
            d = datetime.strptime((a.get("due_date") or "")[:10], "%Y-%m-%d").date()
            delta = (today_dt - d).days
            if delta > oldest_delta:
                oldest_delta = delta
        except Exception:
            pass
    rs.overdue_by_category = by_cat
    rs.overdue_oldest_days = oldest_delta
    rs.overdue_undefined_assignee = undefined_overdue

    # FY26 計画系の期限超過の抜粋（最大8件）
    fy_items = by_cat.get("年度計画策定", [])
    rs.fy_planning_overdue = [
        f"{a.get('due_date','-')} {(a.get('content') or '')[:80]} ({a.get('assignee') or '未定'})"
        for a in fy_items[:8]
    ]

    # 外部依存の抜粋（最大8件）
    ext_items = by_cat.get("外部依存", [])
    rs.external_dep_signals = [
        f"{a.get('due_date','-')} {(a.get('content') or '')[:80]} ({a.get('assignee') or '未定'})"
        for a in ext_items[:8]
    ]

    # 未紐付け率
    period_total = (
        sum(len(m.completed_in_period) + len(m.created_in_period) + len(m.decisions_in_period)
            for m in rd.milestones)
        + len(rd.unlinked_completed) + len(rd.unlinked_created) + len(rd.unlinked_decisions)
    )
    unlinked_total = (
        len(rd.unlinked_completed) + len(rd.unlinked_created) + len(rd.unlinked_decisions)
    )
    rs.no_milestone_count = unlinked_total
    rs.no_milestone_link_ratio = (unlinked_total / period_total) if period_total else 0.0

    # 進捗0%のmilestone
    for m in rd.milestones:
        if m.closed_count == 0 and m.open_count == 0:
            rs.progress_zero_milestones.append(m.milestone_id)
        rem = m.days_remaining()
        if rem is not None and rem <= 30 and (m.closed_count + m.open_count) == 0:
            rs.short_deadline_milestones.append(m.milestone_id)

    # 担当者集中度（期間内に登場する担当者）
    from collections import Counter
    assignee_counter: Counter = Counter()
    undefined_in_period = 0
    for m in rd.milestones:
        for items in (m.completed_in_period, m.created_in_period, m.upcoming_due):
            for a in items:
                v = a.get("assignee")
                if _is_undefined(v):
                    undefined_in_period += 1
                else:
                    # 複数担当はカンマ区切り想定
                    for name in re.split(r"[,、,]", str(v)):
                        n = name.strip()
                        if n:
                            assignee_counter[n] += 1
    rs.assignee_concentration = assignee_counter.most_common(8)
    rs.undefined_assignee_in_period = undefined_in_period

    return rs


def format_risk_signals_for_prompt(rs: RiskSignals) -> str:
    """LLM プロンプトに埋め込む観察事実テキスト。"""
    lines = ["## 観察された機械的シグナル（事実、これらをナラティブで必ず参照すること）"]
    lines.append(f"- 期限超過 AI: {rs.overdue_total} 件 (最古は {rs.overdue_oldest_days} 日経過)")
    cats = rs.overdue_by_category or {}
    lines.append(
        f"  - 年度計画策定系: {len(cats.get('年度計画策定', []))} 件 / "
        f"外部依存系: {len(cats.get('外部依存', []))} 件 / "
        f"その他: {len(cats.get('その他', []))} 件"
    )
    lines.append(f"  - 担当者「未定/null」の超過 AI: {rs.overdue_undefined_assignee} 件")

    if rs.fy_planning_overdue:
        lines.append("- 年度計画策定の期限超過（一部、これらが M1〜M5 すべての前提。最重要）:")
        for s in rs.fy_planning_overdue:
            lines.append(f"  - {s}")

    if rs.external_dep_signals:
        lines.append("- 外部依存（NVIDIA/富士通/契約）の期限超過:")
        for s in rs.external_dep_signals:
            lines.append(f"  - {s}")

    lines.append(f"- 期間内 AI/decision の milestone 未紐付け率: "
                 f"{rs.no_milestone_link_ratio*100:.0f}% ({rs.no_milestone_count} 件)")
    lines.append("  → goals.yaml が現実の業務をカバーできていない可能性、またはタグ付け運用が機能していない")

    if rs.progress_zero_milestones:
        lines.append(f"- 紐付き AI 完了/開 ともに 0 件の milestone: "
                     f"{', '.join(rs.progress_zero_milestones)}")
        lines.append("  → 進捗バー 0% は LLM の主観でしか「進んでいる」と言えない。客観メトリクスが欠如")
    if rs.short_deadline_milestones:
        lines.append(f"- 残 30 日以内かつ紐付き 0 の milestone: "
                     f"{', '.join(rs.short_deadline_milestones)}（クリティカル）")

    if rs.assignee_concentration:
        top = rs.assignee_concentration[:5]
        lines.append("- 担当者集中度（期間内 AI 登場回数 上位5）:")
        for name, n in top:
            lines.append(f"  - {name}: {n} 件")
        lines.append("  → 同一個人への業務集中は属人化リスク。引き継ぎ可能性を評価する必要")
    if rs.undefined_assignee_in_period > 0:
        lines.append(f"- 期間内 AI のうち担当者「未定/null」: {rs.undefined_assignee_in_period} 件")
        lines.append("  → 責任者不在の AI は実行されない。指名が必要")

    return "\n".join(lines)


# --------------------------------------------------------------------------- #
# LLM ナラティブ生成
# --------------------------------------------------------------------------- #
_NARRATIVE_SYSTEM = (
    "あなたは富岳NEXTプロジェクトのベテランPMです。"
    "このレポートは『順調に進捗していることをアピールするもの』ではなく、"
    "意思決定者へのエスカレーションです。どこに障害・リスクがあり、どう対応するかを"
    "明示するために書きます。\n"
    "\n"
    "厳守事項:\n"
    "- 「順調」「概ね順調」「着実に進んでいる」のような楽観的な総括だけで終えない\n"
    "- 「次の2週で追跡」「慎重な対応が求められる」のような曖昧な締めくくりを禁止する\n"
    "- AI/決定/ナレッジの完了件数（処理量）を進捗の根拠にしない\n"
    "  本物の進捗は『成功条件のうちどこが閉じたか』であって、件数ではない\n"
    "- 観察された機械的シグナル（期限超過の年度計画系、外部依存、属人化、紐付け率、"
    "  進捗0%等）が示されている場合、必ずそれに言及し、評価する\n"
    "- 障害には『緩和策』『代替案』『トリガー（いつ判断するか）』『責任者』のいずれかを書く\n"
    "- 担当者が『未定』『null』のAIは責任不明として明記する\n"
    "- 段落の頭に【ステータス】【障害】【対応】【未解決】等の見出しを付けて構造化する\n"
)


def _short(items: list[dict], key: str = "content", n: int = 8) -> str:
    if not items:
        return "(なし)"
    out = []
    for it in items[:n]:
        v = (it.get(key) or "").strip().replace("\n", " ")
        out.append(f"- {v[:160]}")
    return "\n".join(out)


def build_narrative_prompt(m: Milestone, since: str, until: str) -> str:
    parts = [
        f"## マイルストーン {m.milestone_id}: {m.name}",
        f"期限: {m.due_date or '未定'} / 残日数: {m.days_remaining()}",
        f"エリア: {m.area or '-'}",
        f"進捗（紐付き AI ベース）: 完了 {m.closed_count} / オープン {m.open_count}",
        f"対象期間: {since} 〜 {until}",
        "",
    ]
    if m.success_criteria:
        parts.append("### 成功条件（DoDの最重要評価軸）")
        for i, sc in enumerate(m.success_criteria, 1):
            parts.append(f"  {i}. {sc}")
        parts.append("")

    # 機械的に拾った情報
    is_progress_zero = m.closed_count == 0 and m.open_count == 0
    rem = m.days_remaining()
    has_undefined = any(_is_undefined(a.get("assignee")) for a in
                        (m.completed_in_period + m.created_in_period + m.upcoming_due))
    has_external_dep = any(_EXTERNAL_DEP_PAT.search(a.get("content") or "") for a in
                           (m.completed_in_period + m.created_in_period + m.upcoming_due))

    parts.append("### 機械的に観察された注意事項")
    if is_progress_zero:
        parts.append(f"  - 紐付き AI が 0/0。進捗バー 0% で『順調』と判定する根拠は無い")
    if rem is not None and rem <= 60:
        parts.append(f"  - 期限まで残 {rem} 日（短期）")
    if has_undefined:
        parts.append("  - 担当者が『未定/null』の AI が含まれる")
    if has_external_dep:
        parts.append("  - 外部依存（NVIDIA/富士通/契約等）を含む AI が含まれる")
    if not (is_progress_zero or has_undefined or has_external_dep) and (rem is None or rem > 60):
        parts.append("  - （特になし）")
    parts.append("")

    parts.extend([
        f"### 期間内に完了した AI ({len(m.completed_in_period)} 件)",
        _short(m.completed_in_period, n=6),
        "",
        f"### 期間内に新規発生した AI ({len(m.created_in_period)} 件)",
        _short(m.created_in_period, n=5),
        "",
        f"### 期間内の決定事項 ({len(m.decisions_in_period)} 件)",
        _short(m.decisions_in_period, n=5),
        "",
        f"### 次の2週で期限が来る予定 AI ({len(m.upcoming_due)} 件)",
        _short(m.upcoming_due, n=5),
        "",
        "## 出力要件",
        "次の **4段落構成** で書く。各段落の頭に【】見出しを付け、段落間は空行で区切る。",
        "総字数 400〜600字。箇条書き禁止。",
        "",
        "【ステータス】(2-3文)",
        "  進捗判定を 順調 / 注意 / 遅延 / 判断不能 のいずれかで明示する。",
        "  ただし『成功条件のうちどこが閉じ、どこが未着手か』を根拠にして判定する。",
        "  AI 完了件数の多さ・少なさを根拠にしてはいけない（処理量と成果は別物）。",
        "  紐付き AI が 0/0 ならば『判断不能』と書き、客観メトリクス欠如を指摘する。",
        "",
        "【障害・リスク】(3-4文)",
        "  このマイルストーンのゴール達成を阻害する要因を最低 1 件、最大 3 件挙げる。",
        "  外部依存・契約・属人化・期限切迫・前提条件未確定・成功条件のDoD曖昧さなど、",
        "  どの観点でも構わないが、必ず『なぜそれが障害なのか』『放置すると何が起きるか』を書く。",
        "  楽観的な『目立ったリスクなし』で締めることを禁止する。",
        "",
        "【対応・要請】(3-4文)",
        "  上記の障害に対して、次の2週で実行すべきアクションを書く。",
        "  『次の2週で追跡する』のような曖昧な記述ではなく、",
        "  『誰が』『何を』『いつまでに』『どう判断するか』のうち2つ以上を必ず含める。",
        "  PMOや意思決定者への要請（人員・予算・優先度判断）があれば明記する。",
        "",
        "【未解決の論点】(1-2文)",
        "  まだ意思決定されていない事項、合意が取れていない論点を1つ以上挙げる。",
        "  『無し』とする場合は『現時点で意思決定者の判断待ちの論点はないが、",
        "  〇〇については近い将来判断が必要』のように先回り視点を書く。",
        "",
        "出力は4段落の本文のみ（前置き・自己紹介・末尾の総括禁止）。",
    ])
    return "\n".join(parts)


def build_summary_prompt(rd: ReportData) -> str:
    """エグゼクティブサマリー用プロンプト。エスカレーション型で書かせる。"""
    lines = [
        f"# 隔週レポート全体: {rd.since} 〜 {rd.until} (index={rd.index_name})",
        "",
        "## マイルストーン一覧と数値",
    ]
    for m in rd.milestones:
        rem = m.days_remaining()
        rem_s = f"残{rem}日" if rem is not None else "残日数不明"
        lines.append(
            f"- {m.milestone_id} {m.name}: 期限 {m.due_date or '未定'} ({rem_s}), "
            f"紐付きAI 完了/開 {m.closed_count}/{m.open_count}, "
            f"期間内 完了{len(m.completed_in_period)} / 新規{len(m.created_in_period)} / "
            f"決定{len(m.decisions_in_period)} / ナレッジ{len(m.knowledge_in_period)}"
        )
    lines.append("")

    # 機械的シグナル（最重要：このプロンプトの中核）
    lines.append(format_risk_signals_for_prompt(rd.risks))
    lines.append("")

    lines.append("## 各マイルストーンの個別ナラティブ（参考）")
    for m in rd.milestones:
        if m.narrative:
            lines.append(f"### {m.milestone_id}")
            lines.append(m.narrative)
            lines.append("")
    lines.extend([
        "## 出力要件: 1枚もののエグゼクティブサマリー（エスカレーション型）",
        "",
        "このレポートの読み手は意思決定者（部門長・プロジェクトリーダー）であり、",
        "『順調に進捗している』を確認するためのものではない。",
        "**どこに障害があり、どう対応するか、何の判断を求めるか** をエスカレーションする文書である。",
        "",
        "次の **5段落構成** で書く。各段落の頭に【】見出しを付け、段落間は空行で区切る。",
        "総字数 800〜1100字。箇条書き禁止。",
        "",
        "【総合判定】(3-4文)",
        "  プロジェクト全体を 順調 / 注意 / 遅延 / 構造的問題あり のいずれかで判定する。",
        "  前期から状態がどう変わったか（悪化/改善/横這い）を述べる。",
        "  AI 件数の多寡を根拠にせず、成功条件達成の見通しと観察シグナルから判定する。",
        "  紐付き AI 0/0 の milestone が複数ある場合は『進捗計測そのものが機能していない』と書く。",
        "",
        "【今期の主な障害】(4-6文)",
        "  以下の観点を最低3つ取り上げ、それぞれ何が問題で、放置すると何が起きるかを書く:",
        "  (a) 年度計画策定の期限超過があれば、それが M1〜M5 の前提を崩している点",
        "  (b) 外部依存（NVIDIA/富士通/契約）による進行阻害",
        "  (c) 担当者集中・属人化、または『未定/null』担当 AI の責任不在",
        "  (d) 未紐付け AI 比率が高いことの意味（goals.yaml が現実をカバーできていない可能性）",
        "  (e) クリティカルパス上の milestone（特に短期期限）の客観メトリクス欠如",
        "",
        "【今期の前進】(2-3文)",
        "  期間内の意思決定で **構造的に重要** だったものを2-3件だけ簡潔に書く。",
        "  処理量・件数の自慢にならないよう注意。",
        "  『何が決まったから、ゴール達成に近づいたか』を書く。",
        "",
        "【対応・PMOへの要請】(4-5文)",
        "  上記の障害に対し、次の2週でPMOおよび意思決定者が実行/判断すべき事項を書く。",
        "  各項目は『誰が（or どの体制で）』『何を』『いつまでに』『何を判断するか』を含めること。",
        "  曖昧な『追跡する』『慎重に対応する』を禁止。",
        "  必要なら明示的に『〇〇については部門長判断を仰ぐ』とエスカレーションする。",
        "",
        "【未解決のクリティカル論点】(2-3文)",
        "  まだ意思決定されていない、しかし放置できない論点を1〜3件挙げる。",
        "  例: マイルストーン定義の更新、人材配置の見直し、外部協力体制の再交渉、",
        "  クリティカルパス再定義、DoDの明文化等。",
        "  『無し』で締めない。",
        "",
        "出力は5段落の本文のみ（前置き・自己紹介・末尾の総括禁止）。",
    ])
    return "\n".join(lines)


def call_llm_safe(prompt: str, *, system: str = _NARRATIVE_SYSTEM,
                  max_tokens: int = 4096) -> str:
    """LLM 呼び出し + think ブロック除去。

    Kimi-K2-Thinking 系は <think>...</think> を本文より先に長く出すため、
    max_tokens は本文に届く十分な余裕（4096 以上）を確保する。
    閉じタグ未到達で本文が空になった場合は明示メッセージを返す。
    """
    try:
        out = call_argus_llm(prompt, system=system, max_tokens=max_tokens, timeout=240)
        cleaned = strip_think_blocks(out).strip()
        if not cleaned:
            return "(LLM 生成は思考過程のみで本文未到達。max_tokens を増やすか別モデルを試してください)"
        return cleaned
    except Exception as e:
        return f"(LLM 生成に失敗: {e})"


# --------------------------------------------------------------------------- #
# Markdown 出力
# --------------------------------------------------------------------------- #
def render_markdown(rd: ReportData) -> str:
    L = []
    L.append(f"# 隔週進捗レポート ({rd.since} 〜 {rd.until})")
    L.append(f"")
    L.append(f"対象 index: `{rd.index_name}`  /  生成日: {date.today().isoformat()}")
    L.append("")
    L.append("## 全体サマリー")
    L.append("")
    L.append(rd.summary_narrative or "(LLM生成なし)")
    L.append("")

    for m in rd.milestones:
        L.append(f"## {m.milestone_id}: {m.name}")
        L.append(f"- 期限: {m.due_date or '未定'} / 残日数: {m.days_remaining()}")
        L.append(f"- エリア: {m.area or '-'}")
        L.append(f"- 進捗（紐付きAI）: 完了 {m.closed_count} / オープン {m.open_count} = {m.progress_ratio()*100:.0f}%")
        if m.success_criteria:
            L.append("- 成功条件:")
            for sc in m.success_criteria:
                L.append(f"  - {sc}")
        L.append("")
        L.append(f"**ナラティブ**: {m.narrative or '(なし)'}")
        L.append("")
        if m.completed_in_period:
            L.append(f"### 期間内に完了した AI ({len(m.completed_in_period)} 件)")
            for a in m.completed_in_period[:8]:
                L.append(f"- {a.get('content','')[:140]}  — {a.get('assignee') or '担当未定'}")
            L.append("")
        if m.created_in_period:
            L.append(f"### 期間内に新規発生した AI ({len(m.created_in_period)} 件)")
            for a in m.created_in_period[:8]:
                L.append(f"- {a.get('content','')[:140]}  — {a.get('assignee') or '担当未定'}")
            L.append("")
        if m.decisions_in_period:
            L.append(f"### 期間内の決定事項 ({len(m.decisions_in_period)} 件)")
            for d in m.decisions_in_period[:5]:
                L.append(f"- [{d.get('decided_at','')}] {d.get('content','')[:160]}")
            L.append("")
        if m.knowledge_in_period:
            L.append(f"### 期間内に確定したナレッジ ({len(m.knowledge_in_period)} 件)")
            for k in m.knowledge_in_period[:5]:
                L.append(f"- {k['id']}: {k.get('topic','')} — {k.get('current_state','')[:100]}")
            L.append("")
        if m.box_in_period:
            L.append(f"### 期間内に更新された BOX 資料 ({len(m.box_in_period)} 件)")
            for b in m.box_in_period[:8]:
                L.append(f"- [{b.get('modified_at','')}] {b.get('name','')}")
            L.append("")
        if m.upcoming_due:
            L.append(f"### 次の2週で期限が来る予定 AI ({len(m.upcoming_due)} 件)")
            for a in m.upcoming_due[:8]:
                L.append(f"- {a.get('due_date','')} {a.get('content','')[:120]} — {a.get('assignee') or '未定'}")
            L.append("")

    L.append("## 未紐付け（マイルストーン未紐付け）")
    L.append(f"- 期間内完了 AI: {len(rd.unlinked_completed)} 件")
    L.append(f"- 期間内新規 AI: {len(rd.unlinked_created)} 件")
    L.append(f"- 期間内決定事項: {len(rd.unlinked_decisions)} 件")
    L.append("")
    L.append(f"## 期限超過 AI（{len(rd.overdue_items)} 件）")
    for a in rd.overdue_items[:15]:
        L.append(f"- {a.get('due_date','')} {a.get('content','')[:120]} — {a.get('assignee') or '未定'}")
    return "\n".join(L)


# --------------------------------------------------------------------------- #
# pptx 出力
# --------------------------------------------------------------------------- #
def render_pptx(rd: ReportData, out_path: Path) -> None:
    prs = T.make_presentation()
    sw, sh = prs.slide_width, prs.slide_height

    # ----- Slide 1: 表紙 -----
    s = T.blank_slide(prs)
    T.add_bg(s, sw, sh, T.NAVY)
    T.add_rect(s, 0, Inches(5.8), sw, Inches(1.7), T.DEEP)
    T.add_rect(s, 0, Inches(7.3), sw, Inches(0.2), T.TEAL)
    T.add_text(s, Inches(0.8), Inches(1.6), Inches(8), Inches(0.5),
               "BIWEEKLY PROGRESS REPORT", size=14, bold=True, color=T.MINT)
    T.add_rect(s, Inches(0.8), Inches(2.05), Inches(0.8), Inches(0.08), T.CORAL)
    T.add_text(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.5),
               "隔週進捗レポート", size=64, bold=True, color=T.WHITE,
               font=T.HEADER_FONT)
    T.add_text(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.7),
               f"{rd.since} 〜 {rd.until}",
               size=28, color=T.MINT)
    T.add_text(s, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.5),
               f"対象 index: {rd.index_name}", size=14, color=T.MUTED)
    T.add_text(s, Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.3),
               f"生成日 {date.today().isoformat()}  /  富岳NEXT アプリケーション開発エリア",
               size=10, color=T.MUTED)

    # ----- Slide 2: エグゼクティブサマリー（1枚で全体を見渡す） -----
    total_pages = 2 + len(rd.milestones) + 1
    page_no = 2
    s = T.blank_slide(prs)
    T.add_bg(s, sw, sh, T.ICE)
    T.header(s, sh, "E", "エグゼクティブサマリー",
             f"対象期間: {rd.since} 〜 {rd.until}")

    # 左カラム (2/3): 文章中心のナラティブ
    left_x  = Inches(0.55)
    left_y  = Inches(1.95)
    left_w  = Inches(8.4)
    left_h  = Inches(5.20)
    T.add_rect(s, left_x, left_y, left_w, left_h, T.WHITE,
               line=T.MINT, rounded=True, radius=0.02)
    T.add_paragraph_text(s,
                         left_x + Inches(0.15), left_y + Inches(0.10),
                         left_w - Inches(0.3), left_h - Inches(0.2),
                         rd.summary_narrative or "(LLM ナラティブなし)",
                         size=12, color=T.DARK,
                         line_spacing=1.30, paragraph_gap_pt=8,
                         heading_color=T.CORAL)

    # 右カラム (1/3): 主要数値ダッシュボード + マイルストーン一覧
    right_x = Inches(9.10)
    right_y = Inches(1.95)
    right_w = Inches(3.95)

    # KPI カード (3 列) — 「障害シグナル」を主軸に置く
    rs = rd.risks
    fy_n = len(rs.fy_planning_overdue)
    ext_n = len(rs.external_dep_signals)
    overdue_n = rs.overdue_total
    unlinked_pct = int(rs.no_milestone_link_ratio * 100)
    kpi = [
        # 1段目左: 期限超過（赤系）
        (str(overdue_n),   "期限超過 AI",
         T.DANGER if overdue_n >= 20 else (T.CORAL if overdue_n else T.SUCCESS)),
        # 1段目中: FY26 計画策定の遅延（赤）
        (str(fy_n),        "FY26計画 遅延",
         T.DANGER if fy_n > 0 else T.SUCCESS),
        # 1段目右: milestone 未紐付け率
        (f"{unlinked_pct}%",  "未紐付け率",
         T.DANGER if unlinked_pct >= 50 else (T.GOLD if unlinked_pct >= 30 else T.SUCCESS)),
    ]
    kpi_w = (right_w - Inches(0.2)) / 3
    for i, (val, label, clr) in enumerate(kpi):
        kx = right_x + (kpi_w + Inches(0.10)) * i
        T.add_rect(s, kx, right_y, kpi_w, Inches(1.10), T.WHITE,
                   line=T.LINE, rounded=True, radius=0.06)
        T.add_rect(s, kx, right_y, kpi_w, Inches(0.12), clr, rounded=True, radius=0.4)
        T.add_text(s, kx, right_y + Inches(0.18), kpi_w, Inches(0.5),
                   val, size=24, bold=True, color=clr,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        T.add_text(s, kx, right_y + Inches(0.65), kpi_w, Inches(0.45),
                   label, size=9, color=T.GRAY,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # マイルストーン一覧（右カラム下）
    ms_y = right_y + Inches(1.30)
    ms_h_total = left_h - Inches(1.30)
    T.add_rect(s, right_x, ms_y, right_w, ms_h_total, T.WHITE,
               line=T.LINE, rounded=True, radius=0.04)
    T.add_text(s, right_x + Inches(0.15), ms_y + Inches(0.05),
               right_w - Inches(0.3), Inches(0.3),
               "マイルストーン進捗", size=11, bold=True, color=T.TEAL)

    row_top = ms_y + Inches(0.42)
    row_h = (ms_h_total - Inches(0.55)) / max(1, len(rd.milestones))
    for ri, m in enumerate(rd.milestones):
        ry = row_top + row_h * ri
        ratio = m.progress_ratio()
        rem = m.days_remaining()
        rem_s = f"残{rem}日" if rem is not None else "-"
        # 状態色
        if rem is not None and rem < 30 and ratio < 0.7:
            status_clr = T.CORAL
        elif rem is not None and rem < 60 and ratio < 0.5:
            status_clr = T.GOLD
        else:
            status_clr = T.SUCCESS
        # 左サイドの色帯
        T.add_rect(s, right_x + Inches(0.10), ry + Inches(0.05),
                   Inches(0.08), row_h - Inches(0.10),
                   status_clr, rounded=True, radius=0.4)
        # ID + 名称
        T.add_text(s, right_x + Inches(0.25), ry,
                   right_w - Inches(0.4), Inches(0.32),
                   f"{m.milestone_id}  {(m.name or '')[:18]}",
                   size=10, bold=True, color=T.NAVY)
        # 期限・残・期間内数
        info = (
            f"{m.due_date or '未定'} ({rem_s}) · "
            f"完{len(m.completed_in_period)}/新{len(m.created_in_period)}/決{len(m.decisions_in_period)}"
        )
        T.add_text(s, right_x + Inches(0.25), ry + Inches(0.30),
                   right_w - Inches(0.4), Inches(0.28),
                   info, size=8, color=T.GRAY)
        # 進捗バー
        bar_w = right_w - Inches(0.5)
        T.progress_bar(s, right_x + Inches(0.25), ry + Inches(0.58),
                       bar_w, Inches(0.10), ratio)

    T.footer(s, sw, sh, f"対象期間 {rd.since} 〜 {rd.until}", page=page_no, total=total_pages)

    # ----- Milestone slides (文章中心) -----
    for idx, m in enumerate(rd.milestones, start=1):
        page_no += 1
        s = T.blank_slide(prs)
        T.add_bg(s, sw, sh, T.ICE)
        T.header(s, sh, m.milestone_id, m.name,
                 f"期限 {m.due_date or '未定'} / 残 {m.days_remaining()} 日 / "
                 f"完了/開 {m.closed_count}/{m.open_count}")

        # 進捗バー（一行に収める）
        ratio = m.progress_ratio()
        T.add_text(s, Inches(0.55), Inches(1.85), Inches(1.6), Inches(0.32),
                   f"進捗 {ratio*100:.0f}%", size=12, bold=True, color=T.NAVY)
        T.progress_bar(s, Inches(2.0), Inches(1.92), Inches(7.5), Inches(0.20), ratio)
        # 期間内サマリ数値（右端）
        info = (f"期間内: 完{len(m.completed_in_period)} / 新{len(m.created_in_period)} / "
                f"決{len(m.decisions_in_period)} / KN{len(m.knowledge_in_period)}")
        T.add_text(s, Inches(9.7), Inches(1.85), Inches(3.5), Inches(0.32),
                   info, size=10, color=T.GRAY, align=PP_ALIGN.RIGHT)

        # === 左カラム (大): ナラティブ本文 ===
        lx, ly = Inches(0.55), Inches(2.30)
        lw, lh = Inches(8.4), Inches(4.85)
        T.add_rect(s, lx, ly, lw, lh, T.WHITE,
                   line=T.MINT, rounded=True, radius=0.02)
        T.add_text(s, lx + Inches(0.15), ly + Inches(0.08),
                   lw - Inches(0.3), Inches(0.32),
                   "進捗ナラティブ", size=11, bold=True, color=T.TEAL)
        T.add_paragraph_text(s,
                             lx + Inches(0.15), ly + Inches(0.42),
                             lw - Inches(0.3), lh - Inches(0.55),
                             m.narrative or "(LLM ナラティブなし)",
                             size=12, color=T.DARK,
                             line_spacing=1.30, paragraph_gap_pt=6,
                             heading_color=T.CORAL)

        # === 右カラム (小): 成功条件 + 重要 KPI + 次2週の期限 ===
        rx, ry = Inches(9.10), Inches(2.30)
        rw = Inches(3.95)

        # 成功条件
        sc_h = Inches(2.20)
        T.add_rect(s, rx, ry, rw, sc_h, T.WHITE, line=T.LINE, rounded=True, radius=0.04)
        T.add_text(s, rx + Inches(0.12), ry + Inches(0.05),
                   rw - Inches(0.24), Inches(0.30),
                   "成功条件", size=10, bold=True, color=T.TEAL)
        if m.success_criteria:
            sc_items = [sc[:55] for sc in m.success_criteria[:5]]
            T.add_bullets(s,
                          rx + Inches(0.12), ry + Inches(0.36),
                          rw - Inches(0.24), sc_h - Inches(0.42),
                          sc_items, size=9, gap=3)
        else:
            T.add_text(s, rx + Inches(0.12), ry + Inches(0.36),
                       rw - Inches(0.24), Inches(0.4),
                       "(未設定)", size=9, color=T.MUTED, italic=True)

        # 次2週で期限を迎える AI
        nx_y = ry + sc_h + Inches(0.10)
        nx_h = Inches(2.55)
        T.add_rect(s, rx, nx_y, rw, nx_h, T.WHITE, line=T.LINE, rounded=True, radius=0.04)
        T.add_text(s, rx + Inches(0.12), nx_y + Inches(0.05),
                   rw - Inches(0.24), Inches(0.30),
                   f"次2週で期限のAI ({len(m.upcoming_due)}件)",
                   size=10, bold=True, color=T.TEAL)
        if m.upcoming_due:
            items = []
            for a in m.upcoming_due[:6]:
                due = a.get("due_date") or "-"
                content = (a.get("content") or "").replace("\n", " ")[:48]
                assignee = a.get("assignee") or "未定"
                items.append(f"{due} {content} ({assignee})")
            T.add_bullets(s,
                          rx + Inches(0.12), nx_y + Inches(0.36),
                          rw - Inches(0.24), nx_h - Inches(0.42),
                          items, size=9, gap=3)
        else:
            T.add_text(s, rx + Inches(0.12), nx_y + Inches(0.36),
                       rw - Inches(0.24), Inches(0.4),
                       "(該当なし)", size=9, color=T.MUTED, italic=True)

        T.footer(s, sw, sh, f"{m.milestone_id} — {rd.since} 〜 {rd.until}",
                 page=page_no, total=total_pages)

    # ----- 最終スライド: リスクレジスタ（エスカレーション用） -----
    page_no += 1
    s = T.blank_slide(prs)
    T.add_bg(s, sw, sh, T.ICE)
    T.header(s, sh, "R", "リスクレジスタ",
             "意思決定者へのエスカレーション項目")

    rs = rd.risks
    # 上段: FY26 計画策定の遅延（最重要）
    top_x, top_y = Inches(0.55), Inches(1.95)
    top_w, top_h = Inches(12.5), Inches(1.85)
    has_fy = bool(rs.fy_planning_overdue)
    border = T.DANGER if has_fy else T.LINE
    T.add_rect(s, top_x, top_y, top_w, top_h, T.WHITE,
               line=border, rounded=True, radius=0.03)
    T.add_rect(s, top_x, top_y, Inches(0.18), top_h, border, rounded=True, radius=0.4)
    T.add_text(s, top_x + Inches(0.30), top_y + Inches(0.08),
               top_w - Inches(0.5), Inches(0.32),
               f"[最重要] FY2026 年度計画策定の遅延 — {len(rs.fy_planning_overdue)} 件",
               size=12, bold=True, color=T.DANGER if has_fy else T.SUCCESS)
    if has_fy:
        items = rs.fy_planning_overdue[:4]
        T.add_bullets(s, top_x + Inches(0.30), top_y + Inches(0.42),
                      top_w - Inches(0.6), top_h - Inches(0.5),
                      [s_[:140] for s_ in items], size=10, gap=3)
        T.add_text(s, top_x + Inches(0.30), top_y + top_h - Inches(0.33),
                   top_w - Inches(0.6), Inches(0.28),
                   "→ M1〜M5 すべての前提が未確定。責任者の即時指名と5月末までのクローズが必要。",
                   size=10, italic=True, color=T.DANGER)
    else:
        T.add_text(s, top_x + Inches(0.30), top_y + Inches(0.45),
                   top_w - Inches(0.5), Inches(0.4),
                   "FY2026 計画策定の期限超過は検出されていません。",
                   size=10, color=T.MUTED, italic=True)

    # 中段: 左 = 外部依存 / 右 = 属人化
    mid_y = top_y + top_h + Inches(0.12)
    mid_h = Inches(2.45)
    half_w = (top_w - Inches(0.15)) / 2

    # 左: 外部依存
    ext_x = top_x
    has_ext = bool(rs.external_dep_signals)
    border = T.CORAL if has_ext else T.LINE
    T.add_rect(s, ext_x, mid_y, half_w, mid_h, T.WHITE,
               line=border, rounded=True, radius=0.03)
    T.add_rect(s, ext_x, mid_y, Inches(0.15), mid_h, border, rounded=True, radius=0.4)
    T.add_text(s, ext_x + Inches(0.25), mid_y + Inches(0.08),
               half_w - Inches(0.4), Inches(0.32),
               f"外部依存リスク（NVIDIA/富士通/契約） — {len(rs.external_dep_signals)} 件",
               size=11, bold=True, color=T.CORAL if has_ext else T.SUCCESS)
    if has_ext:
        T.add_bullets(s, ext_x + Inches(0.25), mid_y + Inches(0.42),
                      half_w - Inches(0.5), mid_h - Inches(0.55),
                      [s_[:90] for s_ in rs.external_dep_signals[:5]],
                      size=9, gap=3)
    else:
        T.add_text(s, ext_x + Inches(0.25), mid_y + Inches(0.42),
                   half_w - Inches(0.4), Inches(0.4),
                   "(該当なし)", size=10, color=T.MUTED, italic=True)

    # 右: 属人化
    pers_x = ext_x + half_w + Inches(0.15)
    top_assignees = rs.assignee_concentration[:5]
    is_concentrated = bool(top_assignees) and top_assignees[0][1] >= 5
    border = T.GOLD if is_concentrated else T.LINE
    T.add_rect(s, pers_x, mid_y, half_w, mid_h, T.WHITE,
               line=border, rounded=True, radius=0.03)
    T.add_rect(s, pers_x, mid_y, Inches(0.15), mid_h, border, rounded=True, radius=0.4)
    T.add_text(s, pers_x + Inches(0.25), mid_y + Inches(0.08),
               half_w - Inches(0.4), Inches(0.32),
               f"属人化リスク（期間内 AI の集中度）",
               size=11, bold=True, color=T.GOLD if is_concentrated else T.SUCCESS)
    if top_assignees:
        items = [f"{n}: {c} 件" for n, c in top_assignees]
        if rs.undefined_assignee_in_period:
            items.append(f"※ 担当未定/null の AI: {rs.undefined_assignee_in_period} 件")
        T.add_bullets(s, pers_x + Inches(0.25), mid_y + Inches(0.42),
                      half_w - Inches(0.5), mid_h - Inches(0.55),
                      items, size=10, gap=3)

    # 下段: 進捗計測の欠如 + 未紐付け率
    bot_y = mid_y + mid_h + Inches(0.12)
    bot_h = sh - bot_y - Inches(0.50)

    # 左: 進捗計測の欠如
    border = T.DANGER if rs.short_deadline_milestones else T.GOLD
    T.add_rect(s, ext_x, bot_y, half_w, bot_h, T.WHITE,
               line=border, rounded=True, radius=0.03)
    T.add_rect(s, ext_x, bot_y, Inches(0.15), bot_h, border, rounded=True, radius=0.4)
    title_txt = (
        f"進捗計測の欠如（紐付き 0/0 の milestone: {len(rs.progress_zero_milestones)}）"
    )
    T.add_text(s, ext_x + Inches(0.25), bot_y + Inches(0.08),
               half_w - Inches(0.4), Inches(0.32),
               title_txt, size=11, bold=True, color=border)
    msg_lines = []
    if rs.progress_zero_milestones:
        msg_lines.append(
            f"紐付き AI 0/0 の milestone: {', '.join(rs.progress_zero_milestones)}"
        )
    if rs.short_deadline_milestones:
        msg_lines.append(
            f"残30日以内かつ紐付き 0 の milestone: {', '.join(rs.short_deadline_milestones)}"
        )
    msg_lines.append(
        "→ AI を milestone_id で紐付ける運用を即時開始しないと客観進捗は測れない。"
    )
    T.add_bullets(s, ext_x + Inches(0.25), bot_y + Inches(0.42),
                  half_w - Inches(0.5), bot_h - Inches(0.55),
                  msg_lines, size=10, gap=3)

    # 右: 未紐付け率の評価
    pct = int(rs.no_milestone_link_ratio * 100)
    border = T.DANGER if pct >= 50 else (T.GOLD if pct >= 30 else T.LINE)
    T.add_rect(s, pers_x, bot_y, half_w, bot_h, T.WHITE,
               line=border, rounded=True, radius=0.03)
    T.add_rect(s, pers_x, bot_y, Inches(0.15), bot_h, border, rounded=True, radius=0.4)
    T.add_text(s, pers_x + Inches(0.25), bot_y + Inches(0.08),
               half_w - Inches(0.4), Inches(0.32),
               f"goals.yaml カバー率（未紐付け {pct}% = {rs.no_milestone_count} 件）",
               size=11, bold=True, color=border)
    msg_lines = [
        f"期間内 AI/decision のうち {pct}% が M1〜M5 のいずれにも紐付かない。",
    ]
    if pct >= 50:
        msg_lines.append("→ 現実の業務の半分以上が goals.yaml で表現されていない。")
        msg_lines.append("  M6 追加 or 既存 milestone の定義拡張が必要。")
    elif pct >= 30:
        msg_lines.append("→ 紐付けロジック改善 or 新規マイルストーン検討の余地。")
    else:
        msg_lines.append("→ おおむね goals.yaml が現実をカバー。継続観察。")
    T.add_bullets(s, pers_x + Inches(0.25), bot_y + Inches(0.42),
                  half_w - Inches(0.5), bot_h - Inches(0.55),
                  msg_lines, size=10, gap=3)

    T.footer(s, sw, sh, "リスクレジスタ", page=page_no, total=total_pages)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def _draw_block(s, x, y, w, title: str, items: list[str]) -> "Inches":
    box_h = Inches(0.35 + max(1, len(items)) * 0.32 + 0.1)
    T.add_rect(s, x, y, w, box_h, T.WHITE, line=T.LINE, rounded=True, radius=0.04)
    T.add_text(s, x + Inches(0.15), y + Inches(0.05), w - Inches(0.3), Inches(0.3),
               title, size=11, bold=True, color=T.TEAL)
    if items:
        T.add_bullets(s, x + Inches(0.15), y + Inches(0.35),
                      w - Inches(0.3), box_h - Inches(0.4),
                      items, size=10, gap=2)
    else:
        T.add_text(s, x + Inches(0.3), y + Inches(0.35), w - Inches(0.3), Inches(0.3),
                   "(なし)", size=10, color=T.MUTED, italic=True)
    return y + box_h + Inches(0.05)


def _format_ai(items: list[dict], *, n: int = 4, with_due: bool = False) -> list[str]:
    out = []
    for a in items[:n]:
        content = (a.get("content") or "").strip().replace("\n", " ")[:80]
        assignee = a.get("assignee") or "未定"
        if with_due and a.get("due_date"):
            out.append(f"[{a['due_date']}] {content} — {assignee}")
        else:
            out.append(f"{content} — {assignee}")
    return out


def _format_decision(items: list[dict], *, n: int = 3) -> list[str]:
    out = []
    for d in items[:n]:
        content = (d.get("content") or "").strip().replace("\n", " ")[:90]
        out.append(f"[{d.get('decided_at','')}] {content}")
    return out


def _format_knowledge(items: list[dict], *, n: int = 3) -> list[str]:
    out = []
    for k in items[:n]:
        topic = (k.get("topic") or "").strip()[:30]
        cs = (k.get("current_state") or "").strip().replace("\n", " ")[:55]
        out.append(f"{k['id']}: {topic} — {cs}")
    return out


def _format_box(items: list[dict], *, n: int = 4) -> list[str]:
    out = []
    for b in items[:n]:
        name = (b.get("name") or "").replace("\n", " ")[:65]
        date_s = (b.get("modified_at") or "")[:10]
        out.append(f"[{date_s}] {name}")
    return out


# --------------------------------------------------------------------------- #
# メイン
# --------------------------------------------------------------------------- #
def parse_args() -> argparse.Namespace:
    today = date.today()
    default_since = (today - timedelta(days=14)).isoformat()
    default_until = today.isoformat()
    p = argparse.ArgumentParser(description="隔週進捗レポート pptx 生成")
    p.add_argument("--since", default=default_since, help="期間開始 YYYY-MM-DD（既定: 14日前）")
    p.add_argument("--until", default=default_until, help="期間終了 YYYY-MM-DD（既定: 今日）")
    p.add_argument("--index-name", default="pm",
                   help="argus_config.yaml の index 名（既定: pm）")
    p.add_argument("--db", default="data/pm.db", help="pm.db のパス")
    p.add_argument("--knowledge-db", default="data/knowledge.db")
    p.add_argument("--box-db", default="data/box_docs.db")
    p.add_argument("--slack-db", default="data/slack.db")
    p.add_argument("--config", default="data/argus_config.yaml")
    p.add_argument("--output", default=None,
                   help="pptx 出力パス（既定: reports/biweekly_<until>.pptx）")
    p.add_argument("--markdown-only", action="store_true", help="pptx を生成せず .md のみ出力")
    p.add_argument("--md-output", default=None,
                   help="Markdown 出力パス（既定: stdout）")
    p.add_argument("--no-llm", action="store_true", help="LLM ナラティブを生成しない")
    p.add_argument("--lookahead-days", type=int, default=14,
                   help="次N日以内に期限の AI を集める（既定: 14）")
    p.add_argument("--no-encrypt", action="store_true")
    p.add_argument("--dry-run", action="store_true",
                   help="ファイル出力せず標準出力に Markdown を出す")
    add_filter_arg(p)
    return p.parse_args()


def main() -> int:
    args = parse_args()
    cfg = load_argus_config(REPO_ROOT / args.config)
    idx = resolve_index(cfg, args.index_name)

    # filter_presets の解決（--filter が指定された場合）
    filter_channels, filter_meetings = resolve_filter_presets(
        args.filter, config_path=REPO_ROOT / args.config)
    if filter_channels or filter_meetings:
        print(f"[INFO] filter: channels={len(filter_channels)}件, "
              f"meeting_kinds={filter_meetings or '(なし)'}", file=sys.stderr)

    # pm.db
    conn = open_pm_db(REPO_ROOT / args.db, no_encrypt=args.no_encrypt)
    ms_rows = fetch_milestone_progress(conn)
    if not ms_rows:
        print("ERROR: milestones テーブルが空です。pm_ingest.py goals を実行してください。",
              file=sys.stderr)
        return 1

    milestones: list[Milestone] = []
    for r in ms_rows:
        try:
            criteria = json.loads(r.get("success_criteria") or "[]")
        except json.JSONDecodeError:
            criteria = []
        m = Milestone(
            milestone_id=r["milestone_id"],
            goal_id=r.get("goal_id") or "",
            name=r.get("name") or "",
            due_date=r.get("due_date"),
            area=r.get("area"),
            success_criteria=criteria,
            open_count=r.get("open_count") or 0,
            closed_count=r.get("closed_count") or 0,
        )
        m.keywords = extract_keywords(m)
        milestones.append(m)

    # 期間内データ
    completed, created = fetch_period_action_items(conn, args.since, args.until,
                                                   filter_channels, filter_meetings)
    decisions = fetch_period_decisions(conn, args.since, args.until,
                                       filter_channels, filter_meetings)
    upcoming = fetch_upcoming_due(conn, args.until, args.lookahead_days,
                                   filter_channels, filter_meetings)
    overdue = fetch_overdue(conn, args.until, filter_channels, filter_meetings)
    knowledge = fetch_period_knowledge(REPO_ROOT / args.knowledge_db, args.since, args.until)

    # box は index_name 単体で絞る（pm-all は全件相当なので pm-all 除外）
    box_index_filter = None if args.index_name == "pm-all" else [args.index_name]
    box_files = fetch_period_box_files(REPO_ROOT / args.box_db, args.since, args.until,
                                       index_filter=box_index_filter)
    slack_threads = fetch_period_slack_threads(REPO_ROOT / args.slack_db,
                                               args.since, args.until,
                                               idx["channels"])

    # 紐付け
    unlinked_completed, unlinked_created, unlinked_decisions = [], [], []

    def _bin(items: list[dict], target_attr: str, dst_unlinked: list,
             text_key: str = "content"):
        for it in items:
            text = it.get(text_key) or ""
            m = assign_to_milestone(text, milestones, it.get("milestone_id"))
            if m is None:
                dst_unlinked.append(it)
            else:
                getattr(m, target_attr).append(it)

    _bin(completed, "completed_in_period", unlinked_completed)
    _bin(created, "created_in_period", unlinked_created)
    _bin(decisions, "decisions_in_period", unlinked_decisions)
    for k in knowledge:
        text = (k.get("topic") or "") + " " + (k.get("current_state") or "") + " " + (k.get("tags") or "")
        m = assign_to_milestone(text, milestones)
        if m is not None:
            m.knowledge_in_period.append(k)
    for b in box_files:
        text = (b.get("name") or "") + " " + (b.get("folder_path") or "")
        m = assign_to_milestone(text, milestones)
        if m is not None:
            m.box_in_period.append(b)
    for u in upcoming:
        text = u.get("content") or ""
        m = assign_to_milestone(text, milestones, u.get("milestone_id"))
        if m is not None:
            m.upcoming_due.append(u)

    rd = ReportData(
        since=args.since, until=args.until, index_name=args.index_name,
        milestones=milestones,
        unlinked_completed=unlinked_completed,
        unlinked_created=unlinked_created,
        unlinked_decisions=unlinked_decisions,
        overdue_items=overdue,
    )

    # 機械的リスクシグナルを集計（LLM の前に必ず実行）
    rd.risks = collect_risk_signals(rd, args.until)

    # LLM ナラティブ
    if not args.no_llm:
        for m in milestones:
            print(f"[INFO] LLM narrative: {m.milestone_id}", file=sys.stderr)
            m.narrative = call_llm_safe(build_narrative_prompt(m, args.since, args.until),
                                        max_tokens=4096)
        print("[INFO] LLM summary", file=sys.stderr)
        rd.summary_narrative = call_llm_safe(build_summary_prompt(rd), max_tokens=8192)

    # Markdown
    md = render_markdown(rd)
    if args.md_output:
        Path(args.md_output).write_text(md, encoding="utf-8")
        print(f"Markdown 保存: {args.md_output}", file=sys.stderr)
    if args.dry_run or args.markdown_only:
        if not args.md_output:
            print(md)
        if args.markdown_only:
            return 0
        if args.dry_run:
            return 0

    # pptx
    out_path = Path(args.output) if args.output else (
        REPO_ROOT / "reports" / f"biweekly_{args.until}.pptx"
    )
    render_pptx(rd, out_path)
    print(f"pptx 保存: {out_path}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
